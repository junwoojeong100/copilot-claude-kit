from __future__ import annotations

import sys
import unittest
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


SCRIPTS_DIR = Path(__file__).resolve().parents[1] / "scripts"
sys.path.insert(0, str(SCRIPTS_DIR))

from rendered_overlap import (  # noqa: E402
    assign_spans_to_shapes,
    detect_span_overflow,
    detect_span_overlaps,
    shape_records,
)


class RenderedOverlapTests(unittest.TestCase):
    def test_assignment_uses_text_and_nearest_frame(self):
        shapes = [
            {
                "shape": "left",
                "text": "01",
                "normalized": "01",
                "rect": fitz.Rect(0, 0, 50, 50),
            },
            {
                "shape": "right",
                "text": "01",
                "normalized": "01",
                "rect": fitz.Rect(100, 0, 150, 50),
            },
        ]
        spans = [
            {
                "text": "01",
                "normalized": "01",
                "size_pt": 20,
                "rect": fitz.Rect(110, 10, 140, 40),
            }
        ]
        assigned, unmapped = assign_spans_to_shapes(spans, shapes)
        self.assertEqual(unmapped, [])
        self.assertEqual(assigned[0]["shape_index"], 1)

    def test_group_children_and_table_cells_are_individually_mapped(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        group = slide.shapes.add_group_shape()
        child = group.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(0.5)
        )
        child.text = "Grouped text"
        table = slide.shapes.add_table(
            1, 2, Inches(4), Inches(1), Inches(4), Inches(1)
        ).table
        table.cell(0, 0).text = "Left cell"
        table.cell(0, 1).text = "Right cell"
        table.cell(0, 0).merge(table.cell(0, 1))

        document = fitz.open()
        page = document.new_page(width=960, height=540)
        records, unsupported = shape_records(prs, slide, page)
        names = [record["shape"] for record in records]
        self.assertTrue(any("Grouped text" == record["text"] for record in records))
        self.assertTrue(any("cell 1,1" in name for name in names))
        self.assertFalse(any("cell 1,2" in name for name in names))
        merged = next(record for record in records if "cell 1,1" in record["shape"])
        self.assertGreater(merged["rect"].width, 250)
        self.assertEqual(unsupported, [])
        document.close()

    def test_rotated_group_is_classified_for_visual_review(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        group = slide.shapes.add_group_shape()
        textbox = group.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(1)
        )
        textbox.text = "Rotated group"
        group.rotation = 30
        document = fitz.open()
        page = document.new_page(width=960, height=540)

        records, unsupported = shape_records(prs, slide, page)

        self.assertEqual(records, [])
        self.assertEqual(unsupported[0]["kind"], "transformed_group")
        document.close()

    def test_charts_are_explicitly_classified_for_visual_review(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        data = ChartData()
        data.categories = ["A", "B"]
        data.add_series("Series", [1, 2])
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(3),
            data,
        )
        document = fitz.open()
        page = document.new_page(width=960, height=540)
        _, unsupported = shape_records(prs, slide, page)
        self.assertEqual(unsupported[0]["kind"], "chart")
        document.close()

    def test_distinct_text_frames_with_rendered_collision_are_reported(self):
        shapes = [
            {
                "shape": "first",
                "text": "runtime plane",
                "normalized": "runtimeplane",
                "rect": fitz.Rect(0, 0, 100, 30),
            },
            {
                "shape": "second",
                "text": "decision",
                "normalized": "decision",
                "rect": fitz.Rect(70, 25, 150, 55),
            },
        ]
        spans = [
            {
                "shape_index": 0,
                "text": "plane",
                "normalized": "plane",
                "size_pt": 20,
                "rect": fitz.Rect(60, 20, 100, 42),
            },
            {
                "shape_index": 1,
                "text": "decision",
                "normalized": "decision",
                "size_pt": 14,
                "rect": fitz.Rect(80, 30, 140, 48),
            },
        ]
        overlaps = detect_span_overlaps(spans, shapes)
        self.assertEqual(len(overlaps), 1)
        self.assertEqual(overlaps[0]["shape_a"], "first")

        spans[1]["shape_index"] = 0
        self.assertEqual(detect_span_overlaps(spans, shapes), [])

    def test_rendered_spill_is_reported_separately(self):
        shapes = [
            {
                "shape": "wrapped",
                "text": "runtime plane",
                "normalized": "runtimeplane",
                "rect": fitz.Rect(0, 0, 100, 30),
            }
        ]
        spans = [
            {
                "shape_index": 0,
                "text": "runtime",
                "normalized": "runtime",
                "size_pt": 20,
                "rect": fitz.Rect(0, 0, 90, 22),
            },
            {
                "shape_index": 0,
                "text": "plane",
                "normalized": "plane",
                "size_pt": 20,
                "rect": fitz.Rect(20, 24, 65, 48),
            },
        ]
        findings = detect_span_overflow(spans, shapes, tolerance_pt=4)
        self.assertEqual(len(findings), 1)
        self.assertEqual(findings[0]["overflow_edges_pt"]["bottom"], 18)


if __name__ == "__main__":
    unittest.main()
