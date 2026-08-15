from __future__ import annotations

import sys
import tempfile
import unittest
import json
from argparse import Namespace
from pathlib import Path


SCRIPTS_DIR = Path(__file__).resolve().parents[1] / "scripts"
sys.path.insert(0, str(SCRIPTS_DIR))

import render_pptx  # noqa: E402
import language_policy  # noqa: E402
import speaker_notes  # noqa: E402
import toolcheck  # noqa: E402
import tooling  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402
from verify_deck import (  # noqa: E402
    audit_namespace,
    build_parser,
    font_matches,
    internal_fact_id_visibility_failures,
    prepare_output_dirs,
    resolve_contract,
    select_risk_slides,
    source_footer_failures,
    state_label_failures,
    unexpected_fonts,
    verify,
)


class VerifyDeckTests(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(self.temp_dir.cleanup)
        self.work_dir = Path(self.temp_dir.name)

    def test_risk_selection_prioritizes_structural_findings(self):
        report = {
            "text_chars_per_slide": {"values": [100, 500, 200, 300]},
            "small_text_body_candidates": [{"slide": 3}],
            "small_text_label_candidates": [{"slide": 4}],
            "title_risks": [],
            "group_shapes": [],
            "unexpected_out_of_bounds": [{"slide": 1}],
        }
        self.assertEqual(select_risk_slides(report, 3), [1, 3, 2])

    def test_zero_text_decks_do_not_divide_by_zero(self):
        report = {
            "text_chars_per_slide": {"values": [0, 0]},
            "small_text_body_candidates": [],
            "small_text_label_candidates": [],
            "title_risks": [],
            "group_shapes": [],
            "unexpected_out_of_bounds": [],
        }
        self.assertEqual(select_risk_slides(report, 2), [1, 2])

    def test_strict_mode_enables_typography_failures(self):
        args = Namespace(
            deck=self.work_dir / "deck.pptx",
            expected_slides=5,
            allow_bleed="",
            bounds_tolerance=0.02,
            min_body_pt=13.0,
            min_title_pt=26.0,
            title_size_tolerance_pt=0.5,
            footer_top=6.9,
            min_small_text_chars=10,
            fail_small_text=False,
            allow_small_text="",
            allow_overlap="",
            allow_title_size="",
            require_sources="1,3-4",
            fail_unsized_runs=False,
            fail_title_risks=False,
            fail_title_consistency=False,
            fail_overlaps=False,
            strict=True,
        )
        namespace = audit_namespace(args, self.work_dir / "audit.json")
        self.assertTrue(namespace.fail_small_text)
        self.assertTrue(namespace.fail_title_risks)
        self.assertTrue(namespace.fail_title_consistency)
        self.assertTrue(namespace.fail_unsized_runs)
        self.assertTrue(namespace.fail_overlaps)
        self.assertEqual(namespace.require_sources, {1, 3, 4})

        args.allow_small_text = "2"
        namespace = audit_namespace(args, self.work_dir / "audit.json")
        self.assertEqual(namespace.allow_small_text, {2})
        self.assertTrue(namespace.fail_small_text)

        args.allow_overlap = "3"
        namespace = audit_namespace(args, self.work_dir / "audit.json")
        self.assertEqual(namespace.allow_overlap, {3})

        args.allow_title_size = "4"
        namespace = audit_namespace(args, self.work_dir / "audit.json")
        self.assertEqual(namespace.allow_title_size, {4})

    def test_risk_selection_prioritizes_rendered_overlap(self):
        report = {
            "text_chars_per_slide": {"values": [100, 100, 100]},
            "rendered_text_overlaps": [{"slide": 3}],
            "overlap_candidates": [{"slide": 2}],
        }
        self.assertEqual(select_risk_slides(report, 3), [3, 2, 1])

    def test_runner_clears_stale_qa_artifacts(self):
        detail = self.work_dir / "qa-detail"
        render_pptx.claim_output_dir(detail)
        stale = detail / "slide-99.jpg"
        stale.write_text("stale", encoding="utf-8")
        qa_dir, detail_dir = prepare_output_dirs(self.work_dir)
        self.assertTrue(qa_dir.is_dir())
        self.assertTrue(detail_dir.is_dir())
        self.assertFalse(stale.exists())

    def test_runner_refuses_unowned_qa_directory(self):
        qa_dir = self.work_dir / "qa"
        qa_dir.mkdir()
        unrelated = qa_dir / "notes.txt"
        unrelated.write_text("keep", encoding="utf-8")
        with self.assertRaises(RuntimeError):
            prepare_output_dirs(self.work_dir)
        self.assertEqual(unrelated.read_text(encoding="utf-8"), "keep")

    def test_runner_does_not_delete_deck_inside_managed_qa_directory(self):
        out = self.work_dir / "verify"
        qa_dir = out / "qa"
        qa_dir.mkdir(parents=True)
        deck = qa_dir / "deck.pptx"
        deck.write_bytes(b"preserve")
        args = build_parser().parse_args([str(deck), "--out", str(out)])
        with self.assertRaises(ValueError):
            verify(args)
        self.assertEqual(deck.read_bytes(), b"preserve")

    def test_runner_report_cannot_alias_input_deck(self):
        out = self.work_dir / "verify"
        out.mkdir()
        deck = self.work_dir / "deck.pptx"
        deck.write_bytes(b"preserve")
        (out / "verification-report.json").hardlink_to(deck)
        args = build_parser().parse_args([str(deck), "--out", str(out)])
        with self.assertRaises(ValueError):
            verify(args)
        self.assertEqual(deck.read_bytes(), b"preserve")

    def test_font_matching_tolerates_pdf_style_suffixes(self):
        self.assertTrue(
            font_matches(
                "Apple SD Gothic Neo",
                ["ABCDEE+AppleSDGothicNeo-Regular"],
            )
        )
        self.assertFalse(font_matches("Malgun Gothic", ["Aptos-Regular"]))
        self.assertFalse(font_matches("Arial", ["ABCDEF+ArialNarrow-Regular"]))
        self.assertTrue(font_matches("Arial", ["ABCDEF+Arial-BoldItalic"]))
        self.assertEqual(
            unexpected_fonts(
                ["Arial", "Malgun Gothic"],
                ["Arial-Bold", "Aptos", "Malgun Gothic"],
            ),
            ["Aptos"],
        )

    def test_source_footer_uses_human_readable_publishers(self):
        class Context:
            claim_ids_by_slide = {2: ["F-001", "F-002"]}
            fact_ledger = {
                "facts": [
                    {
                        "id": "F-001",
                        "sources": [{"publisher": "Microsoft"}],
                    },
                    {
                        "id": "F-002",
                        "sources": [{"publisher": "GitHub"}],
                    },
                ]
            }

        failures, gaps = source_footer_failures(
            {
                "footer_source_texts_by_slide": {
                    "2": ["출처: Microsoft Learn · 공식 제품 문서"]
                }
            },
            Context(),
        )
        self.assertEqual(gaps, {"2": ["F-002 (GitHub)"]})
        self.assertIn("GitHub", failures[0])

        failures, gaps = source_footer_failures(
            {
                "footer_source_texts_by_slide": {
                    "2": ["출처: Microsoft Learn · GitHub Docs"]
                }
            },
            Context(),
        )
        self.assertEqual(failures, [])
        self.assertEqual(gaps, {})

    def test_internal_fact_ids_are_not_visible(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )
        shape.text = "출처: [F-001] Microsoft Learn"
        deck = self.work_dir / "internal-id.pptx"
        prs.save(deck)

        failures = internal_fact_id_visibility_failures(deck)
        self.assertEqual(len(failures), 1)
        self.assertIn("[F-001]", failures[0])

    def test_without_deck_spec_defaults_remain_backward_compatible(self):
        args = build_parser().parse_args(
            ["deck.pptx", "--out", str(self.work_dir)]
        )
        self.assertIsNone(resolve_contract(args))
        self.assertEqual(args.min_body_pt, 13)
        self.assertEqual(args.min_title_pt, 26)
        self.assertIsNone(args.max_unmapped_text_spans)
        self.assertEqual(args.footer_top, 6.9)

    def test_state_labels_require_standalone_tokens(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(1)
        )
        shape.text = "MEGA platform · PARTIAL GA"
        deck = self.work_dir / "state-label.pptx"
        prs.save(deck)

        class Context:
            spec = {
                "slides": [
                    {
                        "number": 1,
                        "stateLabels": ["GA"],
                    }
                ]
            }

        self.assertTrue(state_label_failures(deck, Context()))

    def test_language_balance_preserves_official_terms_and_flags_english_heavy_slide(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        first = prs.slides.add_slide(prs.slide_layouts[6])
        first.shapes.add_textbox(
            Inches(0.8), Inches(0.8), Inches(11), Inches(1)
        ).text = "GitHub Copilot로 개발 흐름을 개선합니다"
        second = prs.slides.add_slide(prs.slide_layouts[6])
        second.shapes.add_textbox(
            Inches(0.8), Inches(0.8), Inches(11), Inches(1)
        ).text = "English only executive platform operations dashboard controls"
        deck = self.work_dir / "language-balance.pptx"
        prs.save(deck)
        policy = {
            **language_policy.DEFAULT_KOREAN_POLICY,
            "maxLatinRatio": 0.7,
            "maxSlideLatinRatio": 0.6,
            "protectedTerms": ["GitHub Copilot"],
        }
        report = language_policy.analyze_deck(
            deck,
            policy,
            footer_top_in=6.9,
        )
        self.assertEqual(report["missingProtectedTerms"], [])
        self.assertEqual([item["slide"] for item in report["highLatinSlides"]], [2])

    def test_language_balance_requires_configured_official_term(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(
            Inches(0.8), Inches(0.8), Inches(11), Inches(1)
        ).text = "설명 문구만 있습니다"
        deck = self.work_dir / "missing-term.pptx"
        prs.save(deck)
        policy = {
            **language_policy.DEFAULT_KOREAN_POLICY,
            "protectedTerms": ["Microsoft Foundry"],
        }
        report = language_policy.analyze_deck(
            deck,
            policy,
            footer_top_in=6.9,
        )
        self.assertIn("Microsoft Foundry", report["missingProtectedTerms"])

    def test_protected_technical_term_requires_korean_explanation(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(
            Inches(0.8), Inches(0.8), Inches(11), Inches(1)
        ).text = "Hosted Agents"
        deck = self.work_dir / "unexplained-technical-term.pptx"
        prs.save(deck)
        policy = {
            **language_policy.DEFAULT_KOREAN_POLICY,
            "protectedTerms": ["Hosted Agents"],
            "minHangulCharactersPerTechnicalSlide": 12,
        }
        report = language_policy.analyze_deck(
            deck,
            policy,
            footer_top_in=6.9,
        )
        self.assertEqual(
            [item["slide"] for item in report["unexplainedTechnicalSlides"]],
            [1],
        )

    def test_speaker_notes_contract_detects_missing_sections_and_length(self):
        prs = Presentation()
        first = prs.slides.add_slide(prs.slide_layouts[6])
        first.notes_slide.notes_text_frame.text = (
            "질문: 여러분의 회사에서 AI를 쓸수록 실제로 축적되는 것은 무엇입니까?\n"
            "핵심 메시지: 모델보다 기업이 소유한 learning loop가 장기 경쟁력을 만듭니다.\n"
            "전환: 다음 장에서 기업에 남겨야 할 learning asset을 보겠습니다."
        )
        second = prs.slides.add_slide(prs.slide_layouts[6])
        second.notes_slide.notes_text_frame.text = "핵심 메시지: 너무 짧습니다."
        deck = self.work_dir / "speaker-notes.pptx"
        prs.save(deck)
        policy = {
            "required": True,
            "requiredSections": [
                "질문",
                "핵심 메시지",
                "전환",
            ],
            "authoringMode": "regenerate-from-scratch",
            "questionSection": "질문",
            "coreSection": "핵심 메시지",
            "transitionSection": "전환",
            "targetSeconds": 60,
            "minCharacters": 40,
            "maxCharacters": 300,
            "minQuestionCharacters": 20,
            "maxQuestionCharacters": 140,
            "maxQuestionSentences": 1,
            "minCoreCharacters": 30,
            "maxCoreCharacters": 120,
            "maxCoreSentences": 3,
            "maxTransitionCharacters": 100,
            "maxTransitionSentences": 1,
            "maxTotalSentences": 5,
            "requireQuestionFirst": True,
            "requireQuestionMark": True,
            "forbidSourceReferences": True,
        }
        report = speaker_notes.analyze_deck(deck, policy)
        self.assertEqual(report["shortSlides"], [2])
        self.assertEqual(report["sectionGaps"][0]["slide"], 2)
        self.assertIn("질문", report["sectionGaps"][0]["missingSections"])
        self.assertEqual(report["questionNotFirstSlides"], [2])

    def test_speaker_notes_contract_enforces_question_first_and_five_sentences(self):
        prs = Presentation()
        wrong_order = prs.slides.add_slide(prs.slide_layouts[6])
        wrong_order.notes_slide.notes_text_frame.text = (
            "핵심 메시지: 기업이 workflow와 eval을 소유해야 학습이 남습니다.\n"
            "질문: 내일 model을 바꿔도 우리 회사의 기준과 기록은 남습니까?\n"
            "전환: 다음 장에서 운영 loop를 보겠습니다."
        )
        no_question_mark = prs.slides.add_slide(prs.slide_layouts[6])
        no_question_mark.notes_slide.notes_text_frame.text = (
            "질문: 현재 Agent의 개선 근거를 확인할 수 있습니까\n"
            "핵심 메시지: trace와 eval이 있어야 운영 결과를 다음 version으로 되돌릴 수 있습니다.\n"
            "전환: 다음 장에서 promotion gate를 보겠습니다."
        )
        too_many_sentences = prs.slides.add_slide(prs.slide_layouts[6])
        too_many_sentences.notes_slide.notes_text_frame.text = (
            "질문: 첫 pilot의 성공 기준을 한 문장으로 말할 수 있습니까?\n"
            "핵심 메시지: 첫째 문장입니다. 둘째 문장입니다. 셋째 문장입니다. 넷째 문장입니다.\n"
            "전환: 다음 장에서 실행 계획을 보겠습니다."
        )
        source_reference = prs.slides.add_slide(prs.slide_layouts[6])
        source_reference.notes_slide.notes_text_frame.text = (
            "질문: 근거가 발표 흐름보다 더 길어지고 있지는 않습니까?\n"
            "핵심 메시지: 출처는 slide footer와 machine contract에 두고 발표 cue에서는 제거합니다.\n"
            "전환: 다음 장에서 핵심 메시지에 집중하겠습니다.\n"
            "출처: [F-001] https://example.com"
        )
        deck = self.work_dir / "question-first-notes.pptx"
        prs.save(deck)
        policy = {
            "required": True,
            "requiredSections": [
                "질문",
                "핵심 메시지",
                "전환",
            ],
            "authoringMode": "regenerate-from-scratch",
            "questionSection": "질문",
            "coreSection": "핵심 메시지",
            "transitionSection": "전환",
            "targetSeconds": 60,
            "minCharacters": 40,
            "maxCharacters": 500,
            "minQuestionCharacters": 20,
            "maxQuestionCharacters": 140,
            "maxQuestionSentences": 1,
            "minCoreCharacters": 20,
            "maxCoreCharacters": 260,
            "maxCoreSentences": 3,
            "maxTransitionCharacters": 180,
            "maxTransitionSentences": 1,
            "maxTotalSentences": 5,
            "requireQuestionFirst": True,
            "requireQuestionMark": True,
            "forbidSourceReferences": True,
        }
        report = speaker_notes.analyze_deck(deck, policy)
        self.assertEqual(report["questionNotFirstSlides"], [1])
        self.assertEqual(report["questionMarkGaps"], [2])
        self.assertEqual(report["multiSentenceCoreSlides"], [3])
        self.assertEqual(report["overSentenceLimitSlides"], [3])
        self.assertEqual(report["sourceReferenceSlides"], [4])

    def test_speaker_notes_contract_reads_multiline_bullet_sections(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.notes_slide.notes_text_frame.text = (
            "질문:\n"
            "Agent가 무엇을 알고 무엇까지 실행해도 되는지 분리되어 있습니까?\n"
            "핵심 메시지:\n"
            "근거와 실행을 분리해 안전하게 업무를 처리합니다.\n"
            "전환:\n"
            "다음 장에서 운영 통제를 설명합니다."
        )
        deck = self.work_dir / "multiline-notes.pptx"
        prs.save(deck)
        policy = {
            "required": True,
            "requiredSections": [
                "질문",
                "핵심 메시지",
                "전환",
            ],
            "authoringMode": "regenerate-from-scratch",
            "questionSection": "질문",
            "coreSection": "핵심 메시지",
            "transitionSection": "전환",
            "targetSeconds": 60,
            "minCharacters": 40,
            "maxCharacters": 600,
            "minQuestionCharacters": 20,
            "maxQuestionCharacters": 140,
            "maxQuestionSentences": 1,
            "minCoreCharacters": 20,
            "maxCoreCharacters": 100,
            "maxCoreSentences": 3,
            "maxTransitionCharacters": 100,
            "maxTransitionSentences": 1,
            "maxTotalSentences": 5,
            "requireQuestionFirst": True,
            "requireQuestionMark": True,
            "forbidSourceReferences": True,
        }
        report = speaker_notes.analyze_deck(deck, policy)
        self.assertEqual(report["sectionGaps"], [])
        self.assertGreater(report["slides"][0]["questionCharacters"], 20)
        self.assertGreater(report["slides"][0]["coreCharacters"], 20)
        self.assertGreaterEqual(report["slides"][0]["transitionCharacters"], 20)
        self.assertEqual(report["sourceReferenceSlides"], [])

    def test_runner_protects_visual_review_inside_managed_qa(self):
        out = self.work_dir / "verify"
        qa = out / "qa"
        qa.mkdir(parents=True)
        evidence = qa / "visual-review.json"
        evidence.write_text("{}", encoding="utf-8")
        deck = self.work_dir / "deck.pptx"
        deck.write_bytes(b"preserve")
        args = build_parser().parse_args(
            [
                str(deck),
                "--out",
                str(out),
                "--require-visual-review",
                "--visual-review",
                str(evidence),
            ]
        )
        with self.assertRaisesRegex(ValueError, "managed QA"):
            verify(args)
        self.assertTrue(evidence.exists())

    @unittest.skipUnless(
        tooling.resolve_soffice(),
        "LibreOffice is required for verifier integration",
    )
    def test_deck_spec_drives_end_to_end_verification(self):
        font = toolcheck.select_font(
            toolcheck.enumerate_fonts()["fonts"],
            language="en-US",
        )
        self.assertIsNotNone(font)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for text, top, height, size in (
            ("Decision title", 0.6, 0.8, 30),
            ("This body sentence provides enough evidence for the decision.", 2, 0.8, 15),
            ("GA", 5.8, 0.3, 11),
            ("Source: Example · Official documentation", 6.95, 0.2, 8),
        ):
            shape = slide.shapes.add_textbox(
                Inches(0.8), Inches(top), Inches(11.5), Inches(height)
            )
            run = shape.text_frame.paragraphs[0].add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.name = font
        deck = self.work_dir / "contract-deck.pptx"
        prs.save(deck)

        ledger = {
            "schemaVersion": 1,
            "checkedAt": "2026-08-01T10:00:00+09:00",
            "facts": [
                {
                    "id": "F-001",
                    "type": "Fact",
                    "claim": "Official fact",
                    "evidence": "Official documentation",
                    "source": {
                        "title": "Official documentation",
                        "url": "https://example.com/docs",
                    },
                    "publisher": "Example",
                    "publishedOrUpdated": "2026-08-01",
                    "accessed": "2026-08-01",
                    "scopeOrStatus": "GA",
                    "confidence": "High",
                }
            ],
        }
        (self.work_dir / "fact-ledger.json").write_text(
            json.dumps(ledger), encoding="utf-8"
        )
        spec = {
            "schemaVersion": 1,
            "request": {
                "topic": "Decision",
                "audience": "Executive",
                "purpose": "Decision",
                "language": "en-US",
                "slideCount": 1,
            },
            "canvas": {
                "source": "default",
                "widthIn": 13.333,
                "heightIn": 7.5,
            },
            "templateProfile": None,
            "factLedger": "fact-ledger.json",
            "fontPolicy": {
                "selected": font,
                "fallbacks": [],
                "requireAvailable": True,
                "requireRenderedMatch": True,
            },
            "slides": [
                {
                    "number": 1,
                    "role": "evidence",
                    "title": "Decision title",
                    "claimIds": ["F-001"],
                    "stateLabels": ["GA"],
                }
            ],
            "qa": {
                "strict": True,
                "minBodyPt": 15,
                "minTitlePt": 26,
                "maxUnmappedTextSpans": 0,
                "failRenderedOverflow": True,
                "requireVisualReview": False,
                "exceptionManifest": None,
            },
        }
        spec_path = self.work_dir / "deck-spec.json"
        spec_path.write_text(json.dumps(spec), encoding="utf-8")
        args = build_parser().parse_args(
            [
                str(deck),
                "--out",
                str(self.work_dir / "verify"),
                "--deck-spec",
                str(spec_path),
            ]
        )
        result = verify(args)
        self.assertTrue(result["passed"], result["audit_failures"])
        self.assertEqual(result["source_footer_gaps"], {})
        self.assertEqual(result["claim_id_gaps"], {})


if __name__ == "__main__":
    unittest.main()
