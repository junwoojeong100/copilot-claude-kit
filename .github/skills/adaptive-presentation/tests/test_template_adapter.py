from __future__ import annotations

import json
import shutil
import sys
import unittest
from pathlib import Path

SKILL_DIR = Path(__file__).resolve().parents[1]
SCRIPTS_DIR = SKILL_DIR / "scripts"
sys.path.insert(0, str(SKILL_DIR))
sys.path.insert(0, str(SCRIPTS_DIR))

from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from PIL import Image  # noqa: E402

import inspect_template as inspector  # noqa: E402
import pptx_helpers as H  # noqa: E402


class TemplateAdapterTests(unittest.TestCase):
    def setUp(self):
        root = Path(__file__).resolve().parent / ".test-work"
        self.work_dir = root / self._testMethodName
        shutil.rmtree(self.work_dir, ignore_errors=True)
        self.work_dir.mkdir(parents=True)
        self.addCleanup(shutil.rmtree, self.work_dir, True)

    def _template(self) -> Path:
        path = self.work_dir / "template.pptx"
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "SECRET_MARKER"
        prs.save(path)
        return path

    def test_profile_is_deterministic_and_has_contract_fields(self):
        path = self._template()

        first = inspector.inspect_template(path)
        second = inspector.inspect_template(path)

        self.assertEqual(first, second)
        self.assertEqual(first["schemaVersion"], 1)
        self.assertEqual(first["source"], str(path.resolve()))
        self.assertEqual(first["widthIn"], 10.0)
        self.assertEqual(first["heightIn"], 5.625)
        self.assertEqual(first["aspectRatio"], 1.777778)
        self.assertEqual(first["slideCount"], 1)
        self.assertGreaterEqual(first["masterCount"], 1)
        self.assertTrue(first["layouts"])
        self.assertIn("placeholderTypes", first["layouts"][0])
        self.assertEqual(len(first["themeFingerprint"]), 64)
        self.assertEqual(len(first["templateFingerprint"]), 64)
        self.assertIn("latin", first["themeFonts"]["major"])
        self.assertIn("eastAsian", first["themeFonts"]["minor"])
        self.assertNotIn("SECRET_MARKER", json.dumps(first))

    def test_clearing_slides_preserves_template_contract(self):
        path = self._template()
        before = inspector.inspect_template(path)

        prs, _ = H.init_deck(path, clear_existing_slides=True)
        self.assertEqual(len(prs.slides), 0)
        output = self.work_dir / "cleared.pptx"
        prs.save(output)
        after = inspector.inspect_template(output)

        self.assertEqual(after["widthIn"], before["widthIn"])
        self.assertEqual(after["heightIn"], before["heightIn"])
        self.assertEqual(after["masterCount"], before["masterCount"])
        self.assertEqual(after["layouts"], before["layouts"])
        self.assertEqual(after["themeFingerprint"], before["themeFingerprint"])
        self.assertEqual(
            after["templateFingerprint"], before["templateFingerprint"]
        )
        self.assertEqual(after["slideCount"], 0)

    def test_selects_named_or_minimal_placeholder_layout(self):
        path = self._template()

        _, named = H.init_deck(path, layout_name="Title Only")
        _, minimal = H.init_deck(path)

        self.assertEqual(named.name, "Title Only")
        content_types = {
            getattr(placeholder.placeholder_format.type, "name", "")
            for placeholder in minimal.placeholders
        } - {"DATE", "FOOTER", "SLIDE_NUMBER", "HEADER"}
        self.assertEqual(content_types, set())

    def test_slide_authored_media_does_not_change_template_fingerprint(self):
        path = self._template()
        before = inspector.inspect_template(path)
        image = self.work_dir / "slide-image.png"
        Image.new("RGB", (32, 32), "navy").save(image)

        prs, layout = H.init_deck(path, clear_existing_slides=True)
        slide = prs.slides.add_slide(layout)
        slide.shapes.add_picture(str(image), Inches(1), Inches(1))
        output = self.work_dir / "with-slide-media.pptx"
        prs.save(output)

        after = inspector.inspect_template(output)
        self.assertEqual(
            after["templateFingerprint"], before["templateFingerprint"]
        )

    def test_cli_output_refuses_overwrite(self):
        path = self._template()
        output = self.work_dir / "profile.json"

        self.assertEqual(inspector.main([str(path), "--out", str(output)]), 0)
        with self.assertRaises(FileExistsError):
            inspector.main([str(path), "--out", str(output)])


if __name__ == "__main__":
    unittest.main()
