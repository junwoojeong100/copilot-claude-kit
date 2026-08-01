#!/usr/bin/env python3
"""Inspect a PowerPoint template without exposing slide content."""
from __future__ import annotations

import argparse
import hashlib
import json
import posixpath
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree

from pptx import Presentation

EMU_PER_INCH = 914400
SCHEMA_VERSION = 1
DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
RELATIONSHIPS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _placeholder_type_name(placeholder) -> str:
    placeholder_type = placeholder.placeholder_format.type
    name = getattr(placeholder_type, "name", None)
    if name:
        return str(name)
    return str(placeholder_type).split(" (", 1)[0]


def _theme_part(archive: zipfile.ZipFile) -> bytes | None:
    names = sorted(
        name
        for name in archive.namelist()
        if name.startswith("ppt/theme/") and name.endswith(".xml")
    )
    return archive.read(names[0]) if names else None


def _template_fingerprint(archive: zipfile.ZipFile) -> str:
    names = set(archive.namelist())
    pending = [
        name
        for name in names
        if (
            name.startswith("ppt/slideMasters/")
            or name.startswith("ppt/slideLayouts/")
            or name.startswith("ppt/theme/")
        )
        and not name.endswith("/")
        and "/_rels/" not in name
    ]
    included: set[str] = set()
    while pending:
        name = pending.pop()
        if name in included or name not in names:
            continue
        included.add(name)
        directory, filename = posixpath.split(name)
        relationships = f"{directory}/_rels/{filename}.rels"
        if relationships not in names:
            continue
        included.add(relationships)
        try:
            root = ElementTree.fromstring(archive.read(relationships))
        except ElementTree.ParseError:
            continue
        for relationship in root.findall(f"{{{RELATIONSHIPS_NS}}}Relationship"):
            if relationship.attrib.get("TargetMode") == "External":
                continue
            target = relationship.attrib.get("Target")
            if target:
                pending.append(
                    posixpath.normpath(posixpath.join(directory, target)).lstrip("/")
                )

    digest = hashlib.sha256()
    for name in sorted(included):
        digest.update(name.encode("utf-8"))
        digest.update(b"\0")
        digest.update(archive.read(name))
        digest.update(b"\0")
    return digest.hexdigest()


def _typeface(parent, child_name: str) -> str | None:
    if parent is None:
        return None
    child = parent.find(f"{{{DRAWINGML_NS}}}{child_name}")
    if child is None:
        return None
    value = child.attrib.get("typeface", "").strip()
    return value or None


def _theme_fonts(theme_xml: bytes | None) -> dict:
    fonts = {
        "major": {"latin": None, "eastAsian": None},
        "minor": {"latin": None, "eastAsian": None},
    }
    if theme_xml is None:
        return fonts
    try:
        root = ElementTree.fromstring(theme_xml)
    except ElementTree.ParseError:
        return fonts

    scheme = root.find(
        f".//{{{DRAWINGML_NS}}}themeElements/"
        f"{{{DRAWINGML_NS}}}fontScheme"
    )
    if scheme is None:
        return fonts
    for key, element_name in (("major", "majorFont"), ("minor", "minorFont")):
        element = scheme.find(f"{{{DRAWINGML_NS}}}{element_name}")
        fonts[key]["latin"] = _typeface(element, "latin")
        fonts[key]["eastAsian"] = _typeface(element, "ea")
    return fonts


def _layout_inventory(prs: Presentation) -> list[dict]:
    inventory: list[dict] = []
    global_index = 0
    for master_index, master in enumerate(prs.slide_masters):
        for layout_index, layout in enumerate(master.slide_layouts):
            inventory.append(
                {
                    "index": global_index,
                    "masterIndex": master_index,
                    "layoutIndex": layout_index,
                    "name": layout.name or "",
                    "placeholderTypes": [
                        _placeholder_type_name(placeholder)
                        for placeholder in layout.placeholders
                    ],
                }
            )
            global_index += 1
    return inventory


def inspect_template(path: str | Path) -> dict:
    """Return a deterministic, content-safe profile for a PPTX/POTX file."""
    source = Path(path).expanduser().resolve(strict=True)
    if not source.is_file():
        raise ValueError(f"Template path is not a file: {source}")

    prs = Presentation(source)
    with zipfile.ZipFile(source) as archive:
        theme_xml = _theme_part(archive)
        template_fingerprint = _template_fingerprint(archive)

    width_emu = int(prs.slide_width)
    height_emu = int(prs.slide_height)
    return {
        "schemaVersion": SCHEMA_VERSION,
        "source": str(source),
        "widthIn": round(width_emu / EMU_PER_INCH, 6),
        "heightIn": round(height_emu / EMU_PER_INCH, 6),
        "aspectRatio": round(width_emu / height_emu, 6),
        "themeFingerprint": (
            hashlib.sha256(theme_xml).hexdigest() if theme_xml is not None else None
        ),
        "templateFingerprint": template_fingerprint,
        "masterCount": len(prs.slide_masters),
        "layouts": _layout_inventory(prs),
        "themeFonts": _theme_fonts(theme_xml),
        "slideCount": len(prs.slides),
    }


def _json_text(profile: dict) -> str:
    return json.dumps(
        profile,
        ensure_ascii=False,
        indent=2,
        sort_keys=True,
    ) + "\n"


def _write_new_output(source: Path, output: Path, text: str) -> None:
    output = output.expanduser()
    if output.is_symlink() or output.exists():
        raise FileExistsError(f"Refusing to overwrite existing output: {output}")
    if output.resolve(strict=False) == source.resolve(strict=True):
        raise ValueError("Output must not alias the source template")
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("x", encoding="utf-8") as stream:
        stream.write(text)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Emit a deterministic JSON profile for a PPTX/POTX template."
    )
    parser.add_argument("template", type=Path)
    parser.add_argument("--out", type=Path, help="write JSON to a new file")
    args = parser.parse_args(argv)

    profile = inspect_template(args.template)
    text = _json_text(profile)
    if args.out is None:
        sys.stdout.write(text)
    else:
        _write_new_output(Path(profile["source"]), args.out, text)
    return 0


if __name__ == "__main__":
    sys.exit(main())
