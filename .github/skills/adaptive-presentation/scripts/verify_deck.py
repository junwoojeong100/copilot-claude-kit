#!/usr/bin/env python3
"""Run structural audit, full render, risk-slide render, and ZIP checks."""

from __future__ import annotations

import argparse
import json
import re
import shutil
import sys
import zipfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

import audit_pptx
import deck_spec
import inspect_template
import language_policy
import qa_exceptions
import render_pptx
import rendered_overlap
import speaker_notes
import toolcheck
import visual_review
from tooling import path_is_within, paths_collide


def select_risk_slides(report: dict, count: int) -> list[int]:
    values = report["text_chars_per_slide"]["values"]
    maximum = max(max(values, default=0), 1)
    scores = {
        slide: (characters / maximum) * 5
        for slide, characters in enumerate(values, 1)
    }

    def add(items: list[dict], weight: float) -> None:
        for item in items:
            if item.get("slide") is None:
                continue
            slide = int(item["slide"])
            scores[slide] = scores.get(slide, 0) + weight

    add(report.get("small_text_body_candidates", []), 5)
    add(report.get("small_text_label_candidates", []), 1)
    add(report.get("title_risks", []), 10)
    add(report.get("title_size_inconsistencies", []), 12)
    add(report.get("group_shapes", []), 3)
    add(report.get("unsized_runs", []), 5)
    add(report.get("empty_text_frames", []), 1)
    add(report.get("ooxml_repair_risks", []), 30)
    add(report.get("unexpected_out_of_bounds", []), 20)
    add(report.get("overlap_candidates", []), 20)
    add(report.get("rendered_text_overlaps", []), 30)
    add(report.get("rendered_text_overflow_candidates", []), 4)
    add(report.get("unmapped_rendered_text_findings", []), 12)
    add(report.get("unsupported_text_objects", []), 20)
    add(
        report.get("language_balance", {}).get("highLatinSlides", []),
        8,
    )
    add(
        report.get("language_balance", {}).get(
            "unexplainedTechnicalSlides", []
        ),
        10,
    )
    add(
        [
            {"slide": slide}
            for slide in report.get("speaker_notes", {}).get("missingSlides", [])
        ],
        12,
    )
    add(
        [
            {"slide": slide}
            for slide in report.get("speaker_notes", {}).get("shortSlides", [])
        ],
        4,
    )
    add(
        [
            {"slide": item["slide"]}
            for item in report.get("speaker_notes", {}).get("sectionGaps", [])
        ],
        8,
    )
    add(
        [
            {"slide": slide}
            for slide in report.get("speaker_notes", {}).get(
                "briefExplanationSlides", []
            )
        ],
        6,
    )
    add(
        [
            {"slide": slide}
            for slide in report.get("speaker_notes", {}).get(
                "lowExplanationSentenceSlides", []
            )
        ],
        6,
    )
    add(
        [
            {"slide": slide}
            for slide in report.get("missing_required_source_slides", [])
        ],
        15,
    )

    return [
        slide
        for slide, _ in sorted(
            scores.items(),
            key=lambda item: (-item[1], item[0]),
        )[: max(0, count)]
    ]


def parse_slide_list(value: str | None) -> list[int] | None:
    if not value:
        return None
    return render_pptx.parse_slides(value, 10_000)


def audit_namespace(args: argparse.Namespace, report_path: Path) -> argparse.Namespace:
    return argparse.Namespace(
        deck=args.deck,
        expected_slides=args.expected_slides,
        allow_bleed=audit_pptx.parse_slide_set(args.allow_bleed)
        if args.allow_bleed
        else set(),
        bounds_tolerance=args.bounds_tolerance,
        min_body_pt=args.min_body_pt,
        min_title_pt=args.min_title_pt,
        footer_top=args.footer_top,
        min_small_text_chars=args.min_small_text_chars,
        allow_small_text=audit_pptx.parse_slide_set(args.allow_small_text)
        if args.allow_small_text
        else set(),
        allow_overlap=audit_pptx.parse_slide_set(args.allow_overlap)
        if args.allow_overlap
        else set(),
        allow_title_size=audit_pptx.parse_slide_set(args.allow_title_size)
        if args.allow_title_size
        else set(),
        require_sources=audit_pptx.parse_slide_set(args.require_sources)
        if args.require_sources
        else set(),
        title_size_tolerance_pt=args.title_size_tolerance_pt,
        fail_small_text=args.fail_small_text
        or args.strict,
        fail_unsized_runs=args.fail_unsized_runs or args.strict,
        fail_title_risks=args.fail_title_risks or args.strict,
        fail_title_consistency=args.fail_title_consistency or args.strict,
        fail_overlaps=args.fail_overlaps or args.strict,
        json=report_path,
        strict=args.strict,
        allow_finding_ids=getattr(args, "allow_finding_ids", set()),
    )


def render_namespace(
    deck: Path,
    out: Path,
    *,
    reuse_pdf: Path | None = None,
    slides: str | None = None,
    keep_slide_images: bool = False,
) -> argparse.Namespace:
    return argparse.Namespace(
        deck=deck,
        out=out,
        soffice=None,
        conversion_timeout=120.0,
        reuse_pdf=reuse_pdf,
        slides=slides,
        scale=1.25,
        per_sheet=30,
        columns=5,
        thumb_width=220,
        thumb_height=124,
        image_format="jpg",
        quality=82,
        max_image_kb=900,
        keep_slide_images=keep_slide_images,
        keep_pdf=True,
    )


def reconcile(
    name: str,
    cli_value,
    spec_value,
):
    if cli_value is not None and cli_value != spec_value:
        raise ValueError(
            f"{name} conflicts with the authoritative deck spec: "
            f"{cli_value!r} vs {spec_value!r}"
        )
    return spec_value


def resolve_contract(args: argparse.Namespace) -> deck_spec.DeckSpecContext | None:
    context = deck_spec.load_deck_spec(args.deck_spec) if args.deck_spec else None
    if context is None:
        args.min_body_pt = args.min_body_pt or 13.0
        args.min_title_pt = args.min_title_pt or 26.0
        args.footer_top = args.footer_top if args.footer_top is not None else 6.9
        args.fail_rendered_overflow = bool(args.fail_rendered_overflow)
        args.require_visual_review = bool(args.require_visual_review)
        return None

    spec = context.spec
    request = spec["request"]
    qa = spec["qa"]
    args.expected_slides = reconcile(
        "--expected-slides",
        args.expected_slides,
        request["slideCount"],
    )
    args.min_body_pt = reconcile(
        "--min-body-pt",
        args.min_body_pt,
        float(qa["minBodyPt"]),
    )
    args.min_title_pt = reconcile(
        "--min-title-pt",
        args.min_title_pt,
        float(qa["minTitlePt"]),
    )
    args.footer_top = reconcile(
        "--footer-top",
        args.footer_top,
        max(0.0, float(spec["canvas"]["heightIn"]) - 0.6),
    )
    args.max_unmapped_text_spans = reconcile(
        "--max-unmapped-text-spans",
        args.max_unmapped_text_spans,
        int(qa["maxUnmappedTextSpans"]),
    )
    spec_source_slides = context.required_source_slides
    cli_source_slides = (
        audit_pptx.parse_slide_set(args.require_sources)
        if args.require_sources
        else None
    )
    if cli_source_slides is not None and cli_source_slides != spec_source_slides:
        raise ValueError(
            "--require-sources conflicts with claim-bearing slides in the deck spec"
        )
    args.require_sources = ",".join(
        str(slide) for slide in sorted(spec_source_slides)
    )
    args.fail_rendered_overflow = reconcile(
        "--fail-rendered-overflow",
        args.fail_rendered_overflow if args.fail_rendered_overflow else None,
        bool(qa["failRenderedOverflow"]),
    )
    args.require_visual_review = reconcile(
        "--require-visual-review",
        args.require_visual_review if args.require_visual_review else None,
        bool(qa["requireVisualReview"]),
    )
    if qa["strict"]:
        args.strict = True
    return context


def resolve_related_path(
    context: deck_spec.DeckSpecContext | None,
    value: str | None,
) -> Path | None:
    if not value:
        return None
    base = context.path.parent if context is not None else Path.cwd()
    return (base / value).expanduser().resolve()


def normalize_font_name(value: str) -> str:
    name = value.split("+", 1)[-1]
    normalized = "".join(character.casefold() for character in name if character.isalnum())
    suffixes = ("semibold", "regular", "medium", "italic", "bold")
    changed = True
    while changed:
        changed = False
        for suffix in suffixes:
            if normalized.endswith(suffix):
                normalized = normalized[: -len(suffix)]
                changed = True
                break
    normalized = {
        "arialmt": "arial",
    }.get(normalized, normalized)
    return normalized


def font_matches(selected: str, candidates: list[str]) -> bool:
    expected = normalize_font_name(selected)
    return any(
        expected and expected == normalize_font_name(candidate)
        for candidate in candidates
    )


def unexpected_fonts(allowed: list[str], candidates: list[str]) -> list[str]:
    normalized_allowed = {
        normalize_font_name(font)
        for font in allowed
        if normalize_font_name(font)
    }
    return sorted(
        {
            candidate
            for candidate in candidates
            if normalize_font_name(candidate) not in normalized_allowed
        }
    )


def claim_footer_failures(
    audit_report: dict,
    context: deck_spec.DeckSpecContext,
) -> tuple[list[str], dict[str, list[str]]]:
    footers = audit_report.get("footer_source_texts_by_slide", {})
    missing: dict[str, list[str]] = {}
    for slide, claim_ids in context.claim_ids_by_slide.items():
        footer = "\n".join(footers.get(str(slide), []))
        absent = [claim_id for claim_id in claim_ids if f"[{claim_id}]" not in footer]
        if absent:
            missing[str(slide)] = absent
    failures = [
        "Fact Ledger claim IDs missing from source footer on slide "
        f"{slide}: {', '.join(claim_ids)}"
        for slide, claim_ids in missing.items()
    ]
    return failures, missing


def state_label_failures(
    deck: Path,
    context: deck_spec.DeckSpecContext,
) -> list[str]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def visible_text(shape) -> list[str]:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return [
                text
                for child in shape.shapes
                for text in visible_text(child)
            ]
        if getattr(shape, "has_table", False):
            return [
                cell.text
                for row in shape.table.rows
                for cell in row.cells
                if cell.text.strip()
            ]
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            return [shape.text]
        return []

    prs = Presentation(deck)
    state_pattern = re.compile(
        r"(?<![A-Za-z0-9])("
        + "|".join(
            re.escape(label)
            for label in sorted(deck_spec.STATE_LABELS, key=len, reverse=True)
        )
        + r")(?![A-Za-z0-9])",
        re.IGNORECASE,
    )
    failures: list[str] = []
    for slide_spec, slide in zip(context.spec["slides"], prs.slides):
        visible = "\n".join(
            text
            for shape in slide.shapes
            for text in visible_text(shape)
        )
        visible_labels = {
            match.group(1).upper()
            for match in state_pattern.finditer(visible)
        }
        missing = [
            label
            for label in slide_spec["stateLabels"]
            if label.upper() not in visible_labels
        ]
        if missing:
            failures.append(
                f"Required state label(s) missing on slide {slide_spec['number']}: "
                + ", ".join(missing)
            )
    return failures


def all_exception_findings(report: dict) -> list[dict]:
    keys = (
        "out_of_bounds",
        "small_text_body_candidates",
        "small_text_label_candidates",
        "title_risks",
        "title_size_inconsistencies",
        "unsized_runs",
        "overlap_candidates",
        "rendered_text_overlaps",
        "rendered_text_overflow_candidates",
        "unmapped_rendered_text_findings",
        "unsupported_text_objects",
    )
    return [
        finding
        for key in keys
        for finding in report.get(key, [])
        if finding.get("findingId")
    ]


def verify(args: argparse.Namespace) -> dict:
    deck = args.deck.expanduser().resolve()
    contract = resolve_contract(args)
    spec_exception_path = None
    if contract is not None:
        spec_exception_path = contract.spec["qa"].get("exceptionManifest")
    if args.exceptions and spec_exception_path:
        cli_exception_path = Path(args.exceptions).expanduser().resolve()
        contract_exception_path = resolve_related_path(contract, spec_exception_path)
        if cli_exception_path != contract_exception_path:
            raise ValueError(
                "--exceptions conflicts with $.qa.exceptionManifest in the deck spec"
            )
    exception_path = (
        Path(args.exceptions).expanduser().resolve()
        if args.exceptions
        else resolve_related_path(contract, spec_exception_path)
    )
    exception_manifest = qa_exceptions.load_exception_manifest(exception_path)
    args.allow_finding_ids = (
        set(exception_manifest.reasons) if exception_manifest is not None else set()
    )
    out = args.out.expanduser().resolve()
    report_path = out / "verification-report.json"
    contract_inputs = [deck]
    if contract is not None:
        contract_inputs.extend(
            path
            for path in (
                contract.path,
                contract.fact_ledger_path,
                contract.template_profile_path,
            )
            if path is not None
        )
    if exception_path is not None:
        contract_inputs.append(exception_path)
    if args.visual_review is not None:
        contract_inputs.append(args.visual_review.expanduser().resolve())
    for input_path in contract_inputs:
        if paths_collide(input_path, report_path):
            raise ValueError(
                f"Verification report must not overwrite or alias input: {input_path}"
            )
    for candidate in (out / "qa", out / "qa-detail"):
        resolved_candidate = candidate.resolve()
        for input_path in contract_inputs:
            if path_is_within(input_path, resolved_candidate):
                raise ValueError(
                    f"Verification input must be outside managed QA directory "
                    f"{candidate}: {input_path}"
                )
    qa_dir, detail_dir = prepare_output_dirs(out)
    audit_path = qa_dir / "audit.json"

    audit_args = audit_namespace(args, audit_path)
    render_args = render_namespace(deck, qa_dir)
    with ThreadPoolExecutor(max_workers=2) as executor:
        audit_future = executor.submit(audit_pptx.audit, audit_args)
        render_future = executor.submit(render_pptx.render, render_args)
        audit_report, audit_failures = audit_future.result()
        render_manifest = render_future.result()

    with zipfile.ZipFile(deck) as archive:
        corrupt_member = archive.testzip()
    zip_integrity = "ok" if corrupt_member is None else corrupt_member
    if render_manifest["total_slides"] != audit_report["slides"]:
        audit_failures.append(
            "Rendered PDF slide count differs from PPTX: "
            f"{render_manifest['total_slides']} vs {audit_report['slides']}"
        )
    rendered_pdf = render_manifest.get("pdf")
    if not rendered_pdf:
        raise RuntimeError("Verification render did not retain its PDF")
    rendered_findings = rendered_overlap.audit_rendered_text(
        deck,
        Path(rendered_pdf),
        allowed_slides=audit_pptx.parse_slide_set(args.allow_overlap)
        if args.allow_overlap
        else set(),
        allowed_finding_ids=args.allow_finding_ids,
    )
    audit_report.update(rendered_findings)
    if (
        args.fail_overlaps or args.strict
    ) and rendered_findings["unexpected_rendered_text_overlaps"]:
        audit_failures.append(
            f"{len(rendered_findings['unexpected_rendered_text_overlaps'])} "
            "rendered text overlap(s) require repair or --allow-overlap review"
        )
    if args.fail_rendered_overflow and rendered_findings[
        "unexpected_rendered_text_overflow_candidates"
    ]:
        audit_failures.append(
            f"{len(rendered_findings['unexpected_rendered_text_overflow_candidates'])} "
            "rendered text overflow candidate(s) require repair or a finding exception"
        )
    unexpected_unmapped = rendered_findings[
        "unexpected_unmapped_rendered_text_findings"
    ]
    unexpected_unmapped_count = sum(
        finding["count"] for finding in unexpected_unmapped
    )
    if (
        args.max_unmapped_text_spans is not None
        and unexpected_unmapped_count > args.max_unmapped_text_spans
    ):
        audit_failures.append(
            f"{unexpected_unmapped_count} unmapped rendered text span(s) exceed "
            f"the allowed maximum {args.max_unmapped_text_spans}"
        )
    if (
        contract is not None or args.fail_unsupported_text
    ) and rendered_findings["unexpected_unsupported_text_objects"]:
        audit_failures.append(
            f"{len(rendered_findings['unexpected_unsupported_text_objects'])} "
            "unsupported text-bearing object(s) require visual review and a "
            "finding-level exception"
        )

    claim_id_gaps: dict[str, list[str]] = {}
    template_profile_mismatches: list[str] = []
    language_balance = None
    speaker_notes_report = None
    if contract is not None:
        canvas = contract.spec["canvas"]
        actual_width = audit_report["size_inches"]["width"]
        actual_height = audit_report["size_inches"]["height"]
        if (
            abs(actual_width - float(canvas["widthIn"])) > 0.01
            or abs(actual_height - float(canvas["heightIn"])) > 0.01
        ):
            audit_failures.append(
                "Deck canvas differs from deck spec: "
                f"{actual_width}x{actual_height} vs "
                f"{canvas['widthIn']}x{canvas['heightIn']} inches"
            )
        source_failures, claim_id_gaps = claim_footer_failures(
            audit_report, contract
        )
        audit_failures.extend(source_failures)
        audit_failures.extend(state_label_failures(deck, contract))
        if contract.spec.get("languagePolicy") is not None:
            language_balance = language_policy.analyze_deck(
                deck,
                contract.spec["languagePolicy"],
                footer_top_in=args.footer_top,
            )
            audit_report["language_balance"] = language_balance
            audit_failures.extend(language_policy.failures(language_balance))
        if contract.spec.get("speakerNotesPolicy") is not None:
            speaker_notes_report = speaker_notes.analyze_deck(
                deck,
                contract.spec["speakerNotesPolicy"],
            )
            audit_report["speaker_notes"] = speaker_notes_report
            audit_failures.extend(speaker_notes.failures(speaker_notes_report))
        font_policy = contract.spec["fontPolicy"]
        selected_font = font_policy["selected"]
        allowed_fonts = [selected_font, *font_policy["fallbacks"]]
        declared_fonts = [name for name, _ in audit_report["fonts"]]
        if font_policy["requireAvailable"]:
            installed_fonts = toolcheck.enumerate_fonts()["fonts"]
            installed_match = toolcheck.select_font(
                installed_fonts,
                preferred=[selected_font],
                fallbacks=[],
                language=contract.spec["request"]["language"],
            )
            if normalize_font_name(installed_match or "") != normalize_font_name(
                selected_font
            ):
                audit_failures.append(
                    f"Selected font is not installed in the verification environment: "
                    f"{selected_font}"
                )
        if not font_matches(selected_font, declared_fonts):
            audit_failures.append(
                f"Selected font is not explicitly declared in the deck: {selected_font}"
            )
        unexpected_declared = unexpected_fonts(allowed_fonts, declared_fonts)
        if unexpected_declared:
            audit_failures.append(
                "Deck declares fonts outside the selected font policy: "
                + ", ".join(unexpected_declared)
            )
        rendered_fonts = rendered_findings["rendered_fonts"]
        unexpected_rendered = unexpected_fonts(allowed_fonts, rendered_fonts)
        if unexpected_rendered:
            audit_failures.append(
                "Rendered PDF exposes fonts outside the selected font policy: "
                + ", ".join(unexpected_rendered)
            )
        if font_policy["requireRenderedMatch"] and not any(
            font_matches(font, rendered_fonts) for font in allowed_fonts
        ):
            audit_failures.append(
                "Rendered PDF does not expose the selected font or an allowed "
                f"fallback: {selected_font}"
            )
        if contract.template_profile is not None:
            actual_profile = inspect_template.inspect_template(deck)
            expected_profile = contract.template_profile
            for key in (
                "templateFingerprint",
                "themeFingerprint",
                "masterCount",
                "layouts",
            ):
                if actual_profile.get(key) != expected_profile.get(key):
                    template_profile_mismatches.append(key)
            if template_profile_mismatches:
                audit_failures.append(
                    "Deck no longer preserves the template profile field(s): "
                    + ", ".join(template_profile_mismatches)
                )

    used_exception_ids = qa_exceptions.apply_reasons(
        all_exception_findings(audit_report),
        exception_manifest,
    )
    qa_exceptions.reject_stale_exceptions(
        exception_manifest,
        used_exception_ids,
    )
    audit_path.write_text(
        json.dumps(audit_report, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    selected = parse_slide_list(args.risk_slides)
    if selected is None:
        selected = select_risk_slides(audit_report, args.risk_count)
    invalid = [
        slide
        for slide in selected
        if slide < 1 or slide > render_manifest["total_slides"]
    ]
    if invalid:
        raise ValueError(
            f"Risk slides out of range 1-{render_manifest['total_slides']}: {invalid}"
        )

    detail_manifest = None
    if selected:
        pdf = Path(render_manifest["pdf"])
        detail_manifest = render_pptx.render(
            render_namespace(
                deck,
                detail_dir,
                reuse_pdf=pdf,
                slides=",".join(str(slide) for slide in selected),
                keep_slide_images=True,
            )
        )

    automated_passed = not audit_failures and corrupt_member is None
    visual_review_result = None
    visual_review_failure = None
    if args.require_visual_review:
        if args.visual_review is None:
            visual_review_failure = (
                "Visual review evidence is required for this deck revision"
            )
        else:
            try:
                visual_review_result = visual_review.load_visual_review(
                    args.visual_review,
                    deck,
                    slide_count=render_manifest["total_slides"],
                )
            except (
                FileNotFoundError,
                OSError,
                visual_review.VisualReviewError,
            ) as error:
                visual_review_failure = str(error)
    passed = automated_passed and visual_review_failure is None
    result = {
        "deck": str(deck),
        "passed": passed,
        "automated_passed": automated_passed,
        "audit_failures": audit_failures,
        "audit": audit_report,
        "full_render": render_manifest,
        "risk_slides": selected,
        "detail_render": detail_manifest,
        "zip_integrity": zip_integrity,
        "deck_spec": str(contract.path) if contract is not None else None,
        "claim_id_gaps": claim_id_gaps,
        "template_profile_mismatches": template_profile_mismatches,
        "language_balance": language_balance,
        "speaker_notes": speaker_notes_report,
        "exception_manifest": (
            str(exception_manifest.path) if exception_manifest is not None else None
        ),
        "visual_review": (
            {
                "path": str(visual_review_result.path),
                "reviewer": visual_review_result.reviewer,
                "slides": sorted(visual_review_result.reviewed_slides),
            }
            if visual_review_result is not None
            else None
        ),
        "visual_review_failure": visual_review_failure,
    }
    report_path.write_text(
        json.dumps(result, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    result["report"] = str(report_path)
    return result


def prepare_output_dirs(out: Path) -> tuple[Path, Path]:
    qa_dir = out / "qa"
    detail_dir = out / "qa-detail"
    for path in (qa_dir, detail_dir):
        if path.exists():
            if path.is_symlink():
                raise RuntimeError(f"Refusing symlinked QA directory: {path}")
            if not path.is_dir():
                raise RuntimeError(f"Refusing non-directory QA path: {path}")
            if any(path.iterdir()) and not render_pptx.output_dir_is_owned(path):
                raise RuntimeError(
                    f"Refusing to replace non-empty unowned QA directory: {path}"
                )
            shutil.rmtree(path)
        path.mkdir(parents=True, exist_ok=True)
    return qa_dir, detail_dir


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Verify a PPTX with audit, full render, risk slides, and ZIP integrity."
    )
    parser.add_argument("deck", type=Path)
    parser.add_argument("--out", type=Path, required=True)
    parser.add_argument("--deck-spec", type=Path)
    parser.add_argument("--expected-slides", type=audit_pptx.positive_int)
    parser.add_argument("--risk-count", type=render_pptx.positive_int, default=3)
    parser.add_argument("--risk-slides")
    parser.add_argument("--allow-bleed", default="")
    parser.add_argument(
        "--bounds-tolerance", type=audit_pptx.nonnegative_float, default=0.02
    )
    parser.add_argument("--min-body-pt", type=audit_pptx.positive_float)
    parser.add_argument("--min-title-pt", type=audit_pptx.positive_float)
    parser.add_argument("--footer-top", type=audit_pptx.nonnegative_float)
    parser.add_argument(
        "--min-small-text-chars", type=audit_pptx.positive_int, default=10
    )
    parser.add_argument("--fail-small-text", action="store_true")
    parser.add_argument(
        "--allow-small-text",
        default="",
        metavar="SLIDES",
        help="Reviewed slides allowed sub-minimum body text, e.g. 4,8-9.",
    )
    parser.add_argument(
        "--allow-overlap",
        default="",
        metavar="SLIDES",
        help="Reviewed slides allowed intentional geometry/render overlap, e.g. 4,8-9.",
    )
    parser.add_argument(
        "--allow-title-size",
        default="",
        metavar="SLIDES",
        help="Reviewed slides allowed a different content-title size, e.g. 6,12.",
    )
    parser.add_argument(
        "--require-sources",
        default="",
        metavar="SLIDES",
        help="Slides with factual claims that must contain a Source:/출처: footer.",
    )
    parser.add_argument(
        "--title-size-tolerance-pt",
        type=audit_pptx.nonnegative_float,
        default=0.5,
        help="Allowed content-title size variation in points (default: 0.5).",
    )
    parser.add_argument("--fail-unsized-runs", action="store_true")
    parser.add_argument("--fail-title-risks", action="store_true")
    parser.add_argument("--fail-title-consistency", action="store_true")
    parser.add_argument("--fail-overlaps", action="store_true")
    parser.add_argument("--fail-rendered-overflow", action="store_true")
    parser.add_argument("--fail-unsupported-text", action="store_true")
    parser.add_argument(
        "--max-unmapped-text-spans",
        type=audit_pptx.nonnegative_int,
    )
    parser.add_argument(
        "--exceptions",
        type=Path,
        help="Finding-level QA exception manifest with reviewed rationale.",
    )
    parser.add_argument(
        "--require-visual-review",
        action="store_true",
        help="Require revision-bound visual-review evidence before PASS.",
    )
    parser.add_argument("--visual-review", type=Path)
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Enable typography, title, overlap, and configured source failures",
    )
    return parser


def main() -> int:
    args = build_parser().parse_args()
    result = verify(args)
    print(
        f"Verification {'PASS' if result['passed'] else 'FAIL'} | "
        f"slides={result['audit']['slides']} | "
        f"risk={','.join(map(str, result['risk_slides'])) or 'none'}"
    )
    print(f"Report: {result['report']}")
    return 0 if result["passed"] else 1


if __name__ == "__main__":
    sys.exit(main())
