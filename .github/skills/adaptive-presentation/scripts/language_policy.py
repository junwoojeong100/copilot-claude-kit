#!/usr/bin/env python3
"""Korean-first language balance contract and PPTX analysis."""

from __future__ import annotations

import copy
import math
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


HANGUL_RE = re.compile(r"[가-힣]")
LATIN_RE = re.compile(r"[A-Za-z]")
PAGE_NUMBER_RE = re.compile(r"^\d{1,3}$")
SOURCE_PREFIXES = ("source:", "출처:")
STATE_ONLY_RE = re.compile(
    r"^(?:GA|PREVIEW|PARTIAL GA|ASSUMPTION|DEMO DATA)"
    r"(?:\s*[+·|/]\s*(?:GA|PREVIEW|PARTIAL GA|ASSUMPTION|DEMO DATA))*$",
    re.IGNORECASE,
)

DEFAULT_KOREAN_POLICY = {
    "mode": "korean-first-technical-english",
    "targetLatinRatio": 0.40,
    "maxLatinRatio": 0.55,
    "maxSlideLatinRatio": 0.75,
    "minAnalyzedCharacters": 40,
    "preserveOfficialTerms": True,
    "requireKoreanExplanationForProtectedTerms": True,
    "minHangulCharactersPerTechnicalSlide": 24,
    "protectedTerms": [],
    "allowHighLatinSlides": [],
}


class LanguagePolicyError(ValueError):
    """Raised when a language policy is invalid."""


def _ratio(value: Any, path: str) -> float:
    if (
        isinstance(value, bool)
        or not isinstance(value, (int, float))
        or not math.isfinite(value)
        or not 0 <= value <= 1
    ):
        raise LanguagePolicyError(f"{path} must be a finite number between 0 and 1")
    return float(value)


def normalize_policy(
    language: str,
    value: Any,
    *,
    slide_count: int,
) -> dict[str, Any] | None:
    """Return a validated policy, applying Korean defaults when omitted."""
    is_korean = language.casefold() == "ko" or language.casefold().startswith("ko-")
    if value is None:
        return copy.deepcopy(DEFAULT_KOREAN_POLICY) if is_korean else None
    if not isinstance(value, dict):
        raise LanguagePolicyError("$.languagePolicy must be an object or null")

    allowed = set(DEFAULT_KOREAN_POLICY)
    unknown = sorted(set(value) - allowed)
    if unknown:
        raise LanguagePolicyError(
            "$.languagePolicy contains unsupported field(s): " + ", ".join(unknown)
        )
    policy = {**copy.deepcopy(DEFAULT_KOREAN_POLICY), **copy.deepcopy(value)}
    if policy["mode"] != "korean-first-technical-english":
        raise LanguagePolicyError(
            "$.languagePolicy.mode must be 'korean-first-technical-english'"
        )
    if not is_korean:
        raise LanguagePolicyError(
            "$.languagePolicy is only supported when request.language is Korean"
        )
    policy["targetLatinRatio"] = _ratio(
        policy["targetLatinRatio"], "$.languagePolicy.targetLatinRatio"
    )
    policy["maxLatinRatio"] = _ratio(
        policy["maxLatinRatio"], "$.languagePolicy.maxLatinRatio"
    )
    policy["maxSlideLatinRatio"] = _ratio(
        policy["maxSlideLatinRatio"], "$.languagePolicy.maxSlideLatinRatio"
    )
    if policy["maxLatinRatio"] < policy["targetLatinRatio"]:
        raise LanguagePolicyError(
            "$.languagePolicy.maxLatinRatio must be >= targetLatinRatio"
        )
    if policy["maxSlideLatinRatio"] < policy["maxLatinRatio"]:
        raise LanguagePolicyError(
            "$.languagePolicy.maxSlideLatinRatio must be >= maxLatinRatio"
        )
    minimum = policy["minAnalyzedCharacters"]
    if isinstance(minimum, bool) or not isinstance(minimum, int) or minimum < 1:
        raise LanguagePolicyError(
            "$.languagePolicy.minAnalyzedCharacters must be a positive integer"
        )
    if policy["preserveOfficialTerms"] is not True:
        raise LanguagePolicyError(
            "$.languagePolicy.preserveOfficialTerms must be true"
        )
    if policy["requireKoreanExplanationForProtectedTerms"] is not True:
        raise LanguagePolicyError(
            "$.languagePolicy.requireKoreanExplanationForProtectedTerms "
            "must be true"
        )
    minimum_hangul = policy["minHangulCharactersPerTechnicalSlide"]
    if (
        isinstance(minimum_hangul, bool)
        or not isinstance(minimum_hangul, int)
        or minimum_hangul < 1
    ):
        raise LanguagePolicyError(
            "$.languagePolicy.minHangulCharactersPerTechnicalSlide "
            "must be a positive integer"
        )
    terms = policy["protectedTerms"]
    if (
        not isinstance(terms, list)
        or not all(isinstance(term, str) and term.strip() for term in terms)
        or len(terms) != len(set(terms))
    ):
        raise LanguagePolicyError(
            "$.languagePolicy.protectedTerms must contain unique non-empty strings"
        )
    policy["protectedTerms"] = [term.strip() for term in terms]
    allowed_slides = policy["allowHighLatinSlides"]
    if (
        not isinstance(allowed_slides, list)
        or not all(
            isinstance(slide, int)
            and not isinstance(slide, bool)
            and 1 <= slide <= slide_count
            for slide in allowed_slides
        )
        or len(allowed_slides) != len(set(allowed_slides))
    ):
        raise LanguagePolicyError(
            "$.languagePolicy.allowHighLatinSlides must contain unique valid slide numbers"
        )
    return policy


def _shape_texts(shape, *, inherited_top: int | None = None):
    top = inherited_top if inherited_top is not None else getattr(shape, "top", None)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from _shape_texts(child, inherited_top=top)
        return
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    yield top, cell.text
        return
    if getattr(shape, "has_text_frame", False) and shape.text.strip():
        yield top, shape.text


def _analyzable_text(text: str) -> str:
    lines = []
    for line in text.splitlines():
        stripped = " ".join(line.split())
        lowered = stripped.casefold()
        if not stripped or PAGE_NUMBER_RE.fullmatch(stripped):
            continue
        if any(lowered.startswith(prefix) for prefix in SOURCE_PREFIXES):
            continue
        if STATE_ONLY_RE.fullmatch(stripped):
            continue
        lines.append(stripped)
    return " ".join(lines)


def _counts(text: str) -> tuple[int, int, float]:
    latin = len(LATIN_RE.findall(text))
    hangul = len(HANGUL_RE.findall(text))
    total = latin + hangul
    return latin, hangul, latin / total if total else 0.0


def _neutralize_terms(text: str, terms: list[str]) -> str:
    neutralized = text
    for term in sorted(terms, key=len, reverse=True):
        neutralized = re.sub(
            re.escape(term),
            " ",
            neutralized,
            flags=re.IGNORECASE,
        )
    return neutralized


def analyze_deck(
    deck: Path,
    policy: dict[str, Any],
    *,
    footer_top_in: float,
) -> dict[str, Any]:
    """Measure visible non-footer Korean/Latin copy and protected terms."""
    prs = Presentation(deck)
    footer_top = int(footer_top_in * 914400)
    slides: list[dict[str, Any]] = []
    all_text: list[str] = []
    unexplained_technical_slides: list[dict[str, Any]] = []
    for number, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            for top, text in _shape_texts(shape):
                if top is not None and top >= footer_top:
                    continue
                cleaned = _analyzable_text(text)
                if cleaned:
                    parts.append(cleaned)
        visible = " ".join(parts)
        all_text.append(visible)
        raw_latin, raw_hangul, raw_ratio = _counts(visible)
        matched_terms = [
            term
            for term in policy["protectedTerms"]
            if term.casefold() in visible.casefold()
        ]
        latin, hangul, ratio = _counts(
            _neutralize_terms(visible, policy["protectedTerms"])
        )
        slides.append(
            {
                "slide": number,
                "rawLatinCharacters": raw_latin,
                "rawHangulCharacters": raw_hangul,
                "rawLatinRatio": round(raw_ratio, 4),
                "protectedTerms": matched_terms,
                "latinCharacters": latin,
                "hangulCharacters": hangul,
                "analyzedCharacters": latin + hangul,
                "latinRatio": round(ratio, 4),
            }
        )
        if (
            policy["requireKoreanExplanationForProtectedTerms"]
            and matched_terms
            and raw_hangul < policy["minHangulCharactersPerTechnicalSlide"]
        ):
            unexplained_technical_slides.append(
                {
                    "slide": number,
                    "hangulCharacters": raw_hangul,
                    "protectedTerms": matched_terms,
                }
            )
    combined = " ".join(all_text)
    raw_latin, raw_hangul, raw_ratio = _counts(combined)
    latin, hangul, ratio = _counts(
        _neutralize_terms(combined, policy["protectedTerms"])
    )
    allowed = set(policy["allowHighLatinSlides"])
    high_slides = [
        item
        for item in slides
        if item["analyzedCharacters"] >= policy["minAnalyzedCharacters"]
        and item["latinRatio"] > policy["maxSlideLatinRatio"]
        and item["slide"] not in allowed
    ]
    folded = combined.casefold()
    missing_terms = [
        term
        for term in policy["protectedTerms"]
        if term.casefold() not in folded
    ]
    return {
        "mode": policy["mode"],
        "targetLatinRatio": policy["targetLatinRatio"],
        "maxLatinRatio": policy["maxLatinRatio"],
        "maxSlideLatinRatio": policy["maxSlideLatinRatio"],
        "overall": {
            "rawLatinCharacters": raw_latin,
            "rawHangulCharacters": raw_hangul,
            "rawLatinRatio": round(raw_ratio, 4),
            "latinCharacters": latin,
            "hangulCharacters": hangul,
            "analyzedCharacters": latin + hangul,
            "latinRatio": round(ratio, 4),
        },
        "slides": slides,
        "highLatinSlides": high_slides,
        "protectedTerms": policy["protectedTerms"],
        "missingProtectedTerms": missing_terms,
        "unexplainedTechnicalSlides": unexplained_technical_slides,
    }


def failures(report: dict[str, Any]) -> list[str]:
    problems: list[str] = []
    overall = report["overall"]
    if (
        overall["analyzedCharacters"]
        and overall["latinRatio"] > report["maxLatinRatio"]
    ):
        problems.append(
            "Deck Latin-character ratio exceeds Korean-first policy: "
            "(protected official terms excluded) "
            f"{overall['latinRatio']:.1%} > {report['maxLatinRatio']:.1%} "
            f"(target {report['targetLatinRatio']:.1%})"
        )
    if report["highLatinSlides"]:
        slides = ", ".join(str(item["slide"]) for item in report["highLatinSlides"])
        problems.append(
            "Slide Latin-character ratio exceeds the configured maximum on slide(s): "
            + slides
        )
    if report["missingProtectedTerms"]:
        problems.append(
            "Protected official service/feature term(s) are missing from the deck: "
            + ", ".join(report["missingProtectedTerms"])
        )
    if report["unexplainedTechnicalSlides"]:
        rendered = "; ".join(
            f"{item['slide']} ({', '.join(item['protectedTerms'])})"
            for item in report["unexplainedTechnicalSlides"]
        )
        problems.append(
            "Slides containing protected technical terms need more simple "
            "Korean explanation: " + rendered
        )
    return problems
