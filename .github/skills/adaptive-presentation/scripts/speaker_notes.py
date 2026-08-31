#!/usr/bin/env python3
"""Speaker-notes contract and PPTX analysis."""

from __future__ import annotations

import copy
import re
from pathlib import Path
from typing import Any

from pptx import Presentation


CORE_ONLY_MODE = "core-only"
GUIDED_FLOW_MODE = "guided-flow"


DEFAULT_KOREAN_POLICY = {
    "mode": CORE_ONLY_MODE,
    "required": True,
    "authoringMode": "regenerate-from-scratch",
    "requiredSections": ["핵심 메시지"],
    "coreSection": "핵심 메시지",
    "targetSeconds": 60,
    "minCharacters": 120,
    "maxCharacters": 600,
    "minCoreCharacters": 100,
    "maxCoreCharacters": 580,
    "minCoreSentences": 4,
    "maxCoreSentences": 6,
    "minTotalSentences": 4,
    "maxTotalSentences": 6,
    "forbiddenSections": ["질문", "전환"],
    "requireCoreFirst": True,
    "forbidSourceReferences": True,
}


GUIDED_FLOW_POLICY = {
    "mode": GUIDED_FLOW_MODE,
    "required": True,
    "authoringMode": "regenerate-from-scratch",
    "requiredSections": ["질문", "핵심 메시지", "전환"],
    "questionSection": "질문",
    "coreSection": "핵심 메시지",
    "transitionSection": "전환",
    "targetSeconds": 60,
    "minCharacters": 80,
    "maxCharacters": 600,
    "minQuestionCharacters": 20,
    "maxQuestionCharacters": 140,
    "maxQuestionSentences": 1,
    "minCoreCharacters": 30,
    "maxCoreCharacters": 260,
    "minCoreSentences": 1,
    "maxCoreSentences": 3,
    "maxTransitionCharacters": 180,
    "maxTransitionSentences": 1,
    "minTotalSentences": 1,
    "maxTotalSentences": 5,
    "requireQuestionFirst": True,
    "requireQuestionMark": True,
    "forbidSourceReferences": True,
}


class SpeakerNotesPolicyError(ValueError):
    """Raised when a speaker-notes policy is invalid."""


def _infer_mode(value: dict[str, Any]) -> str:
    explicit = value.get("mode")
    if explicit is not None:
        if explicit not in {CORE_ONLY_MODE, GUIDED_FLOW_MODE}:
            raise SpeakerNotesPolicyError(
                "$.speakerNotesPolicy.mode must be core-only or guided-flow"
            )
        return explicit
    guided_keys = {
        "questionSection",
        "transitionSection",
        "minQuestionCharacters",
        "maxQuestionCharacters",
        "maxQuestionSentences",
        "maxTransitionCharacters",
        "maxTransitionSentences",
        "requireQuestionFirst",
        "requireQuestionMark",
    }
    sections = value.get("requiredSections")
    if any(key in value for key in guided_keys) or (
        isinstance(sections, list)
        and any(section in {"질문", "전환"} for section in sections)
    ):
        return GUIDED_FLOW_MODE
    return CORE_ONLY_MODE


def _validate_sections(policy: dict[str, Any]) -> None:
    sections = policy["requiredSections"]
    if (
        not isinstance(sections, list)
        or not all(isinstance(section, str) and section.strip() for section in sections)
        or len(sections) != len(set(sections))
    ):
        raise SpeakerNotesPolicyError(
            "$.speakerNotesPolicy.requiredSections must contain unique "
            "non-empty strings"
        )
    policy["requiredSections"] = [section.strip() for section in sections]
    section_keys = ["coreSection"]
    if policy["mode"] == GUIDED_FLOW_MODE:
        section_keys.extend(["questionSection", "transitionSection"])
    for key in section_keys:
        section = policy[key]
        if (
            not isinstance(section, str)
            or section.strip() not in policy["requiredSections"]
        ):
            raise SpeakerNotesPolicyError(
                f"$.speakerNotesPolicy.{key} must name one of requiredSections"
            )
        policy[key] = section.strip()
    if policy["mode"] == CORE_ONLY_MODE:
        forbidden = policy["forbiddenSections"]
        if (
            not isinstance(forbidden, list)
            or not all(
                isinstance(section, str) and section.strip()
                for section in forbidden
            )
            or len(forbidden) != len(set(forbidden))
        ):
            raise SpeakerNotesPolicyError(
                "$.speakerNotesPolicy.forbiddenSections must contain unique "
                "non-empty strings"
            )
        policy["forbiddenSections"] = [
            section.strip() for section in forbidden
        ]
        overlap = sorted(
            set(policy["requiredSections"]) & set(policy["forbiddenSections"])
        )
        if overlap:
            raise SpeakerNotesPolicyError(
                "$.speakerNotesPolicy sections cannot be both required and "
                "forbidden: " + ", ".join(overlap)
            )


def normalize_policy(language: str, value: Any) -> dict[str, Any] | None:
    """Return a validated policy, applying Korean defaults when omitted."""
    is_korean = language.casefold() == "ko" or language.casefold().startswith("ko-")
    if value is None:
        return copy.deepcopy(DEFAULT_KOREAN_POLICY) if is_korean else None
    if not isinstance(value, dict):
        raise SpeakerNotesPolicyError(
            "$.speakerNotesPolicy must be an object or null"
        )
    mode = _infer_mode(value)
    defaults = (
        DEFAULT_KOREAN_POLICY
        if mode == CORE_ONLY_MODE
        else GUIDED_FLOW_POLICY
    )
    allowed = set(defaults)
    unknown = sorted(set(value) - allowed)
    if unknown:
        raise SpeakerNotesPolicyError(
            "$.speakerNotesPolicy contains unsupported field(s): "
            + ", ".join(unknown)
        )
    policy = {**copy.deepcopy(defaults), **copy.deepcopy(value), "mode": mode}
    if not isinstance(policy["required"], bool):
        raise SpeakerNotesPolicyError(
            "$.speakerNotesPolicy.required must be a boolean"
        )
    if policy["authoringMode"] != "regenerate-from-scratch":
        raise SpeakerNotesPolicyError(
            "$.speakerNotesPolicy.authoringMode must be "
            "'regenerate-from-scratch'"
        )
    _validate_sections(policy)
    boolean_keys = ["forbidSourceReferences"]
    if mode == CORE_ONLY_MODE:
        boolean_keys.append("requireCoreFirst")
    else:
        boolean_keys.extend(["requireQuestionFirst", "requireQuestionMark"])
    for key in boolean_keys:
        if not isinstance(policy[key], bool):
            raise SpeakerNotesPolicyError(
                f"$.speakerNotesPolicy.{key} must be a boolean"
            )
    integer_keys = [
        "minCharacters",
        "maxCharacters",
        "targetSeconds",
        "minCoreCharacters",
        "maxCoreCharacters",
        "minCoreSentences",
        "maxCoreSentences",
        "minTotalSentences",
        "maxTotalSentences",
    ]
    if mode == GUIDED_FLOW_MODE:
        integer_keys.extend(
            [
                "minQuestionCharacters",
                "maxQuestionCharacters",
                "maxQuestionSentences",
                "maxTransitionCharacters",
                "maxTransitionSentences",
            ]
        )
    for key in integer_keys:
        item = policy[key]
        if isinstance(item, bool) or not isinstance(item, int) or item < 1:
            raise SpeakerNotesPolicyError(
                f"$.speakerNotesPolicy.{key} must be a positive integer"
            )
    if policy["maxCharacters"] < policy["minCharacters"]:
        raise SpeakerNotesPolicyError(
            "$.speakerNotesPolicy.maxCharacters must be >= minCharacters"
        )
    if (
        mode == GUIDED_FLOW_MODE
        and policy["maxQuestionCharacters"] < policy["minQuestionCharacters"]
    ):
        raise SpeakerNotesPolicyError(
            "$.speakerNotesPolicy.maxQuestionCharacters must be >= "
            "minQuestionCharacters"
        )
    if policy["maxCoreCharacters"] < policy["minCoreCharacters"]:
        raise SpeakerNotesPolicyError(
            "$.speakerNotesPolicy.maxCoreCharacters must be >= "
            "minCoreCharacters"
        )
    if policy["maxCoreSentences"] < policy["minCoreSentences"]:
        raise SpeakerNotesPolicyError(
            "$.speakerNotesPolicy.maxCoreSentences must be >= "
            "minCoreSentences"
        )
    if policy["maxTotalSentences"] < policy["minTotalSentences"]:
        raise SpeakerNotesPolicyError(
            "$.speakerNotesPolicy.maxTotalSentences must be >= "
            "minTotalSentences"
        )
    return policy


def _notes_text(slide) -> str:
    try:
        value = slide.notes_slide.notes_text_frame.text
    except (AttributeError, ValueError):
        return ""
    return value.strip()


def _parse_sections(text: str, sections: list[str]) -> dict[str, str]:
    parsed = {section: [] for section in sections}
    ordered_sections = sorted(sections, key=len, reverse=True)
    current: str | None = None
    for line in text.splitlines():
        stripped = line.strip()
        matched = next(
            (
                section
                for section in ordered_sections
                if stripped.startswith(f"{section}:")
            ),
            None,
        )
        if matched is not None:
            current = matched
            remainder = stripped[len(matched) + 1 :].strip()
            if remainder:
                parsed[current].append(remainder)
            continue
        if current is not None and stripped:
            parsed[current].append(stripped)
    return {
        section: "\n".join(lines).strip()
        for section, lines in parsed.items()
    }


def _sentence_count(text: str) -> int:
    return len(re.findall(r"[.!?](?:\s|$)", text))


def _contains_source_reference(text: str) -> bool:
    return any(
        re.search(pattern, text, flags)
        for pattern, flags in (
            (r"^\s*(?:출처|source)\s*:", re.IGNORECASE | re.MULTILINE),
            (r"\[?F-\d{3,}\]?", re.IGNORECASE),
            (r"https?://", re.IGNORECASE),
        )
    )


def _present_section_labels(text: str, sections: list[str]) -> list[str]:
    return [
        section
        for section in sections
        if re.search(
            rf"^\s*{re.escape(section)}\s*:",
            text,
            re.MULTILINE,
        )
    ]


def analyze_deck(
    deck: Path,
    policy: dict[str, Any],
) -> dict[str, Any]:
    """Inspect notes length and required section markers on every slide."""
    normalized = normalize_policy("ko-KR", policy)
    if normalized is None:
        raise SpeakerNotesPolicyError("Speaker notes policy is required")
    policy = normalized
    mode = policy["mode"]
    prs = Presentation(deck)
    slides: list[dict[str, Any]] = []
    missing_slides: list[int] = []
    short_slides: list[int] = []
    long_slides: list[int] = []
    section_gaps: list[dict[str, Any]] = []
    core_not_first_slides: list[int] = []
    forbidden_section_slides: list[dict[str, Any]] = []
    question_not_first_slides: list[int] = []
    brief_question_slides: list[int] = []
    long_question_slides: list[int] = []
    multi_sentence_question_slides: list[int] = []
    question_mark_gaps: list[int] = []
    brief_core_slides: list[int] = []
    long_core_slides: list[int] = []
    under_core_sentence_minimum_slides: list[int] = []
    multi_sentence_core_slides: list[int] = []
    long_transition_slides: list[int] = []
    multi_sentence_transition_slides: list[int] = []
    under_sentence_minimum_slides: list[int] = []
    over_sentence_limit_slides: list[int] = []
    source_reference_slides: list[int] = []
    for number, slide in enumerate(prs.slides, 1):
        text = _notes_text(slide)
        character_count = len(text)
        total_sentences = _sentence_count(text)
        has_source_reference = _contains_source_reference(text)
        parsed_section_names = list(policy["requiredSections"])
        if mode == CORE_ONLY_MODE:
            parsed_section_names.extend(
                section
                for section in policy["forbiddenSections"]
                if section not in parsed_section_names
            )
        parsed_sections = _parse_sections(text, parsed_section_names)
        missing_sections = [
            section
            for section in policy["requiredSections"]
            if not parsed_sections[section]
        ]
        if not text:
            missing_slides.append(number)
        elif character_count < policy["minCharacters"]:
            short_slides.append(number)
        if character_count > policy["maxCharacters"]:
            long_slides.append(number)
        if missing_sections:
            section_gaps.append(
                {
                    "slide": number,
                    "missingSections": missing_sections,
                }
            )
        core = parsed_sections[policy["coreSection"]]
        core_characters = len(core)
        core_sentences = _sentence_count(core)
        core_first = text.startswith(f"{policy['coreSection']}:")
        if (
            mode == CORE_ONLY_MODE
            and text
            and policy["requireCoreFirst"]
            and not core_first
        ):
            core_not_first_slides.append(number)
        if core and core_characters < policy["minCoreCharacters"]:
            brief_core_slides.append(number)
        if core_characters > policy["maxCoreCharacters"]:
            long_core_slides.append(number)
        if core and core_sentences < policy["minCoreSentences"]:
            under_core_sentence_minimum_slides.append(number)
        if core_sentences > policy["maxCoreSentences"]:
            multi_sentence_core_slides.append(number)

        question_characters = 0
        question_sentences = 0
        question_first = False
        has_question_mark = False
        transition_characters = 0
        transition_sentences = 0
        forbidden_sections_present: list[str] = []
        if mode == GUIDED_FLOW_MODE:
            question = parsed_sections[policy["questionSection"]]
            question_characters = len(question)
            question_sentences = _sentence_count(question)
            question_first = text.startswith(
                f"{policy['questionSection']}:"
            )
            has_question_mark = question.rstrip().endswith(("?", "？"))
            if text and policy["requireQuestionFirst"] and not question_first:
                question_not_first_slides.append(number)
            if (
                question
                and question_characters < policy["minQuestionCharacters"]
            ):
                brief_question_slides.append(number)
            if question_characters > policy["maxQuestionCharacters"]:
                long_question_slides.append(number)
            if question_sentences > policy["maxQuestionSentences"]:
                multi_sentence_question_slides.append(number)
            if (
                question
                and policy["requireQuestionMark"]
                and not has_question_mark
            ):
                question_mark_gaps.append(number)
            transition = parsed_sections[policy["transitionSection"]]
            transition_characters = len(transition)
            transition_sentences = _sentence_count(transition)
            if transition_characters > policy["maxTransitionCharacters"]:
                long_transition_slides.append(number)
            if transition_sentences > policy["maxTransitionSentences"]:
                multi_sentence_transition_slides.append(number)
        else:
            forbidden_sections_present = _present_section_labels(
                text,
                policy["forbiddenSections"],
            )
            if forbidden_sections_present:
                forbidden_section_slides.append(
                    {
                        "slide": number,
                        "sections": forbidden_sections_present,
                    }
                )

        if text and total_sentences < policy["minTotalSentences"]:
            under_sentence_minimum_slides.append(number)
        if total_sentences > policy["maxTotalSentences"]:
            over_sentence_limit_slides.append(number)
        if policy["forbidSourceReferences"] and has_source_reference:
            source_reference_slides.append(number)
        slides.append(
            {
                "slide": number,
                "characters": character_count,
                "sentences": total_sentences,
                "missingSections": missing_sections,
                "coreFirst": core_first,
                "forbiddenSectionsPresent": forbidden_sections_present,
                "questionCharacters": question_characters,
                "questionSentences": question_sentences,
                "questionFirst": question_first,
                "hasQuestionMark": has_question_mark,
                "coreCharacters": core_characters,
                "coreSentences": core_sentences,
                "transitionCharacters": transition_characters,
                "transitionSentences": transition_sentences,
                "hasSourceReference": has_source_reference,
            }
        )
    return {
        "mode": mode,
        "required": policy["required"],
        "authoringMode": policy["authoringMode"],
        "requiredSections": policy["requiredSections"],
        "questionSection": policy.get("questionSection"),
        "coreSection": policy["coreSection"],
        "transitionSection": policy.get("transitionSection"),
        "forbiddenSections": policy.get("forbiddenSections", []),
        "minCharacters": policy["minCharacters"],
        "maxCharacters": policy["maxCharacters"],
        "targetSeconds": policy["targetSeconds"],
        "minQuestionCharacters": policy.get("minQuestionCharacters"),
        "maxQuestionCharacters": policy.get("maxQuestionCharacters"),
        "maxQuestionSentences": policy.get("maxQuestionSentences"),
        "minCoreCharacters": policy["minCoreCharacters"],
        "maxCoreCharacters": policy["maxCoreCharacters"],
        "minCoreSentences": policy["minCoreSentences"],
        "maxCoreSentences": policy["maxCoreSentences"],
        "maxTransitionCharacters": policy.get("maxTransitionCharacters"),
        "maxTransitionSentences": policy.get("maxTransitionSentences"),
        "minTotalSentences": policy["minTotalSentences"],
        "maxTotalSentences": policy["maxTotalSentences"],
        "requireCoreFirst": policy.get("requireCoreFirst", False),
        "requireQuestionFirst": policy.get("requireQuestionFirst", False),
        "requireQuestionMark": policy.get("requireQuestionMark", False),
        "forbidSourceReferences": policy["forbidSourceReferences"],
        "slides": slides,
        "missingSlides": missing_slides,
        "shortSlides": short_slides,
        "longSlides": long_slides,
        "sectionGaps": section_gaps,
        "coreNotFirstSlides": core_not_first_slides,
        "forbiddenSectionSlides": forbidden_section_slides,
        "questionNotFirstSlides": question_not_first_slides,
        "briefQuestionSlides": brief_question_slides,
        "longQuestionSlides": long_question_slides,
        "multiSentenceQuestionSlides": multi_sentence_question_slides,
        "questionMarkGaps": question_mark_gaps,
        "briefCoreSlides": brief_core_slides,
        "longCoreSlides": long_core_slides,
        "underCoreSentenceMinimumSlides": (
            under_core_sentence_minimum_slides
        ),
        "multiSentenceCoreSlides": multi_sentence_core_slides,
        "longTransitionSlides": long_transition_slides,
        "multiSentenceTransitionSlides": multi_sentence_transition_slides,
        "underSentenceMinimumSlides": under_sentence_minimum_slides,
        "overSentenceLimitSlides": over_sentence_limit_slides,
        "sourceReferenceSlides": source_reference_slides,
    }


def failures(report: dict[str, Any]) -> list[str]:
    if not report["required"]:
        return []
    problems: list[str] = []
    if report["missingSlides"]:
        problems.append(
            "Speaker notes are missing on slide(s): "
            + ", ".join(map(str, report["missingSlides"]))
        )
    if report["shortSlides"]:
        problems.append(
            "Speaker notes are shorter than "
            f"{report['minCharacters']} characters on slide(s): "
            + ", ".join(map(str, report["shortSlides"]))
        )
    if report["longSlides"]:
        problems.append(
            "Speaker notes exceed "
            f"{report['maxCharacters']} characters on slide(s): "
            + ", ".join(map(str, report["longSlides"]))
        )
    if report["sectionGaps"]:
        rendered = "; ".join(
            f"{item['slide']} ({', '.join(item['missingSections'])})"
            for item in report["sectionGaps"]
        )
        problems.append(
            "Speaker notes are missing required section(s): " + rendered
        )
    if report["coreNotFirstSlides"]:
        problems.append(
            "Speaker notes do not start with the core-message section on "
            "slide(s): "
            + ", ".join(map(str, report["coreNotFirstSlides"]))
        )
    if report["forbiddenSectionSlides"]:
        rendered = "; ".join(
            f"{item['slide']} ({', '.join(item['sections'])})"
            for item in report["forbiddenSectionSlides"]
        )
        problems.append(
            "Speaker notes contain forbidden section(s): " + rendered
        )
    if report["questionNotFirstSlides"]:
        problems.append(
            "Speaker notes do not start with the question section on slide(s): "
            + ", ".join(map(str, report["questionNotFirstSlides"]))
        )
    if report["briefQuestionSlides"]:
        problems.append(
            "Speaker note questions are shorter than "
            f"{report['minQuestionCharacters']} characters on slide(s): "
            + ", ".join(map(str, report["briefQuestionSlides"]))
        )
    if report["longQuestionSlides"]:
        problems.append(
            "Speaker note questions exceed "
            f"{report['maxQuestionCharacters']} characters on slide(s): "
            + ", ".join(map(str, report["longQuestionSlides"]))
        )
    if report["multiSentenceQuestionSlides"]:
        problems.append(
            "Speaker note questions contain more than "
            f"{report['maxQuestionSentences']} sentence(s) on slide(s): "
            + ", ".join(map(str, report["multiSentenceQuestionSlides"]))
        )
    if report["questionMarkGaps"]:
        problems.append(
            "Speaker note questions do not end with a question mark on slide(s): "
            + ", ".join(map(str, report["questionMarkGaps"]))
        )
    if report["briefCoreSlides"]:
        problems.append(
            "Speaker note core messages are shorter than "
            f"{report['minCoreCharacters']} characters on slide(s): "
            + ", ".join(map(str, report["briefCoreSlides"]))
        )
    if report["longCoreSlides"]:
        problems.append(
            "Speaker note core messages exceed "
            f"{report['maxCoreCharacters']} characters on slide(s): "
            + ", ".join(map(str, report["longCoreSlides"]))
        )
    if report["underCoreSentenceMinimumSlides"]:
        problems.append(
            "Speaker note core messages contain fewer than "
            f"{report['minCoreSentences']} sentence(s) on slide(s): "
            + ", ".join(
                map(str, report["underCoreSentenceMinimumSlides"])
            )
        )
    if report["multiSentenceCoreSlides"]:
        problems.append(
            "Speaker note core messages contain more than "
            f"{report['maxCoreSentences']} sentence(s) on slide(s): "
            + ", ".join(map(str, report["multiSentenceCoreSlides"]))
        )
    if report["longTransitionSlides"]:
        problems.append(
            "Speaker note transitions exceed "
            f"{report['maxTransitionCharacters']} characters on slide(s): "
            + ", ".join(map(str, report["longTransitionSlides"]))
        )
    if report["multiSentenceTransitionSlides"]:
        problems.append(
            "Speaker note transitions contain more than "
            f"{report['maxTransitionSentences']} sentence(s) on slide(s): "
            + ", ".join(map(str, report["multiSentenceTransitionSlides"]))
        )
    if report["underSentenceMinimumSlides"]:
        problems.append(
            "Speaker notes contain fewer than "
            f"{report['minTotalSentences']} total sentence(s) on slide(s): "
            + ", ".join(map(str, report["underSentenceMinimumSlides"]))
        )
    if report["overSentenceLimitSlides"]:
        problems.append(
            "Speaker notes exceed "
            f"{report['maxTotalSentences']} total sentences on slide(s): "
            + ", ".join(map(str, report["overSentenceLimitSlides"]))
        )
    if report["sourceReferenceSlides"]:
        problems.append(
            "Speaker notes contain forbidden source references on slide(s): "
            + ", ".join(map(str, report["sourceReferenceSlides"]))
        )
    return problems
