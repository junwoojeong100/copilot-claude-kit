#!/usr/bin/env python3
"""Detect rendered text spill and collisions in a PPTX-derived PDF."""

from __future__ import annotations

import math
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from qa_exceptions import annotate


def normalize_text(text: str) -> str:
    return "".join(character.casefold() for character in text if character.isalnum())


def distance_to_rect(rect: fitz.Rect, x: float, y: float) -> float:
    dx = max(rect.x0 - x, 0.0, x - rect.x1)
    dy = max(rect.y0 - y, 0.0, y - rect.y1)
    return math.hypot(dx, dy)


def page_rect(
    prs: Presentation,
    page: fitz.Page,
    left: float,
    top: float,
    width: float,
    height: float,
) -> fitz.Rect:
    x_scale = page.rect.width / prs.slide_width
    y_scale = page.rect.height / prs.slide_height
    return fitz.Rect(
        page.rect.x0 + left * x_scale,
        page.rect.y0 + top * y_scale,
        page.rect.x0 + (left + width) * x_scale,
        page.rect.y0 + (top + height) * y_scale,
    )


def group_transform(
    shape,
    parent_transform: tuple[float, float, float, float, float, float] | None,
) -> tuple[float, float, float, float, float, float]:
    xfrm = shape._element.grpSpPr.xfrm
    child_width = max(float(xfrm.chExt.cx), 1.0)
    child_height = max(float(xfrm.chExt.cy), 1.0)
    left, top, width, height = absolute_bounds(shape, parent_transform)
    return (
        left,
        top,
        width / child_width,
        height / child_height,
        float(xfrm.chOff.x),
        float(xfrm.chOff.y),
    )


def group_has_unsupported_transform(shape) -> bool:
    xfrm = shape._element.grpSpPr.xfrm
    rotation = int(xfrm.get("rot") or 0)
    flip_h = str(xfrm.get("flipH") or "").casefold() in {"1", "true"}
    flip_v = str(xfrm.get("flipV") or "").casefold() in {"1", "true"}
    return rotation % 21600000 != 0 or flip_h or flip_v


def absolute_bounds(
    shape,
    transform: tuple[float, float, float, float, float, float] | None,
) -> tuple[float, float, float, float]:
    if transform is None:
        return float(shape.left), float(shape.top), float(shape.width), float(shape.height)
    origin_x, origin_y, scale_x, scale_y, offset_x, offset_y = transform
    return (
        origin_x + (float(shape.left) - offset_x) * scale_x,
        origin_y + (float(shape.top) - offset_y) * scale_y,
        float(shape.width) * scale_x,
        float(shape.height) * scale_y,
    )


def table_records(
    prs: Presentation,
    shape,
    page: fitz.Page,
    bounds: tuple[float, float, float, float],
    z_index: int,
    name_prefix: str,
) -> list[dict]:
    left, top, width, height = bounds
    table_width = max(sum(float(column.width) for column in shape.table.columns), 1.0)
    table_height = max(sum(float(row.height) for row in shape.table.rows), 1.0)
    width_scale = width / table_width
    height_scale = height / table_height
    records: list[dict] = []
    row_top = top
    for row_index, row in enumerate(shape.table.rows):
        column_left = left
        for column_index, column in enumerate(shape.table.columns):
            cell = shape.table.cell(row_index, column_index)
            if cell.is_spanned:
                column_left += float(column.width) * width_scale
                continue
            text = " ".join(cell.text.split())
            if text:
                normalized = normalize_text(text)
                cell_width = sum(
                    float(shape.table.columns[index].width)
                    for index in range(
                        column_index,
                        min(
                            column_index + int(cell.span_width),
                            len(shape.table.columns),
                        ),
                    )
                )
                cell_height = sum(
                    float(shape.table.rows[index].height)
                    for index in range(
                        row_index,
                        min(
                            row_index + int(cell.span_height),
                            len(shape.table.rows),
                        ),
                    )
                )
                records.append(
                    {
                        "z_index": z_index,
                        "shape": (
                            f"{name_prefix} cell {row_index + 1},{column_index + 1}"
                        ),
                        "shape_id": shape.shape_id,
                        "text": text,
                        "normalized": normalized,
                        "rect": page_rect(
                            prs,
                            page,
                            column_left,
                            row_top,
                            cell_width * width_scale,
                            cell_height * height_scale,
                        ),
                        "kind": "table_cell",
                    }
                )
            column_left += float(column.width) * width_scale
        row_top += float(row.height) * height_scale
    return records


def shape_records(prs: Presentation, slide, page: fitz.Page) -> tuple[list[dict], list[dict]]:
    records: list[dict] = []
    unsupported: list[dict] = []

    def visit(
        shapes,
        *,
        transform: tuple[float, float, float, float, float, float] | None = None,
        prefix: str = "",
    ) -> None:
        for z_index, shape in enumerate(shapes):
            shape_name = getattr(shape, "name", f"shape-{shape.shape_id}")
            qualified_name = f"{prefix}/{shape_name}" if prefix else shape_name
            bounds = absolute_bounds(shape, transform)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                if group_has_unsupported_transform(shape):
                    unsupported.append(
                        {
                            "kind": "transformed_group",
                            "shape": qualified_name,
                            "reason": (
                                "Rotated or flipped group text mapping requires "
                                "visual review."
                            ),
                        }
                    )
                    continue
                visit(
                    shape.shapes,
                    transform=group_transform(shape, transform),
                    prefix=qualified_name,
                )
                continue
            if getattr(shape, "has_chart", False):
                unsupported.append(
                    {
                        "kind": "chart",
                        "shape": qualified_name,
                        "reason": "Chart-rendered labels require visual review.",
                    }
                )
                continue
            if (
                shape.shape_type
                in {
                    MSO_SHAPE_TYPE.DIAGRAM,
                    MSO_SHAPE_TYPE.IGX_GRAPHIC,
                }
            ):
                unsupported.append(
                    {
                        "kind": "graphic_frame",
                        "shape": qualified_name,
                        "reason": "SmartArt or diagram text mapping is not supported.",
                    }
                )
                continue
            if getattr(shape, "has_table", False):
                records.extend(
                    table_records(
                        prs,
                        shape,
                        page,
                        bounds,
                        z_index,
                        qualified_name,
                    )
                )
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = " ".join(shape.text.split())
            normalized = normalize_text(text)
            if not normalized:
                continue
            left, top, width, height = bounds
            records.append(
                {
                    "z_index": z_index,
                    "shape": qualified_name,
                    "shape_id": shape.shape_id,
                    "text": text,
                    "normalized": normalized,
                    "rect": page_rect(prs, page, left, top, width, height),
                    "kind": "text_frame",
                }
            )

    visit(slide.shapes)
    return records, unsupported


def span_records(page: fitz.Page) -> list[dict]:
    records: list[dict] = []
    for block in page.get_text("dict")["blocks"]:
        if block.get("type") != 0:
            continue
        for line in block["lines"]:
            for span in line["spans"]:
                text = " ".join(span["text"].split())
                normalized = normalize_text(text)
                if not normalized:
                    continue
                records.append(
                    {
                        "text": text,
                        "normalized": normalized,
                        "size_pt": round(float(span["size"]), 2),
                        "font": str(span.get("font") or ""),
                        "rect": fitz.Rect(span["bbox"]),
                    }
                )
    return records


def assign_spans_to_shapes(
    spans: list[dict],
    shapes: list[dict],
) -> tuple[list[dict], list[dict]]:
    assigned: list[dict] = []
    unmapped: list[dict] = []
    for span in spans:
        rect = span["rect"]
        center_x = (rect.x0 + rect.x1) / 2
        center_y = (rect.y0 + rect.y1) / 2
        candidates: list[tuple[float, int]] = []
        for shape_index, shape in enumerate(shapes):
            span_text = span["normalized"]
            shape_text = shape["normalized"]
            if span_text not in shape_text and shape_text not in span_text:
                continue
            frame = shape["rect"]
            frame_center_x = (frame.x0 + frame.x1) / 2
            frame_center_y = (frame.y0 + frame.y1) / 2
            score = distance_to_rect(frame, center_x, center_y) * 10
            score += math.hypot(
                center_x - frame_center_x,
                center_y - frame_center_y,
            ) / max(len(span_text), 1)
            candidates.append((score, shape_index))
        if not candidates:
            unmapped.append(span)
            continue
        record = dict(span)
        record["shape_index"] = min(candidates)[1]
        assigned.append(record)
    return assigned, unmapped


def rect_values(rect: fitz.Rect) -> list[float]:
    return [round(value, 2) for value in (rect.x0, rect.y0, rect.x1, rect.y1)]


def detect_span_overlaps(
    assigned_spans: list[dict],
    shapes: list[dict],
    *,
    min_ratio: float = 0.03,
    min_dimension_pt: float = 1.0,
) -> list[dict]:
    overlaps: dict[tuple[int, int], dict] = {}
    for index, first in enumerate(assigned_spans):
        for second in assigned_spans[index + 1 :]:
            first_shape = int(first["shape_index"])
            second_shape = int(second["shape_index"])
            if first_shape == second_shape:
                continue
            intersection = first["rect"] & second["rect"]
            if (
                intersection.is_empty
                or intersection.width < min_dimension_pt
                or intersection.height < min_dimension_pt
            ):
                continue
            smaller_area = min(first["rect"].get_area(), second["rect"].get_area())
            ratio = intersection.get_area() / max(smaller_area, 1.0)
            if ratio < min_ratio:
                continue
            key = tuple(sorted((first_shape, second_shape)))
            item = {
                "shape_a": shapes[first_shape]["shape"],
                "shape_a_text": shapes[first_shape]["text"][:160],
                "span_a": first["text"][:120],
                "shape_b": shapes[second_shape]["shape"],
                "shape_b_text": shapes[second_shape]["text"][:160],
                "span_b": second["text"][:120],
                "overlap_ratio": round(ratio, 3),
                "overlap_rect_pt": rect_values(intersection),
            }
            previous = overlaps.get(key)
            if previous is None or item["overlap_ratio"] > previous["overlap_ratio"]:
                overlaps[key] = item
    return list(overlaps.values())


def detect_span_overflow(
    assigned_spans: list[dict],
    shapes: list[dict],
    *,
    tolerance_pt: float = 4.0,
) -> list[dict]:
    by_shape: dict[int, list[dict]] = {}
    for span in assigned_spans:
        by_shape.setdefault(int(span["shape_index"]), []).append(span)

    findings: list[dict] = []
    for shape_index, spans in by_shape.items():
        rendered = fitz.Rect(spans[0]["rect"])
        for span in spans[1:]:
            rendered |= span["rect"]
        frame = shapes[shape_index]["rect"]
        edges = {
            "left": max(0.0, frame.x0 - rendered.x0),
            "top": max(0.0, frame.y0 - rendered.y0),
            "right": max(0.0, rendered.x1 - frame.x1),
            "bottom": max(0.0, rendered.y1 - frame.y1),
        }
        worst = max(edges.values())
        if worst <= tolerance_pt:
            continue
        findings.append(
            {
                "shape": shapes[shape_index]["shape"],
                "text": shapes[shape_index]["text"][:160],
                "max_overflow_pt": round(worst, 2),
                "overflow_edges_pt": {
                    edge: round(value, 2)
                    for edge, value in edges.items()
                    if value > tolerance_pt
                },
                "frame_rect_pt": rect_values(frame),
                "rendered_rect_pt": rect_values(rendered),
            }
        )
    return findings


def audit_rendered_text(
    deck: Path,
    pdf: Path,
    *,
    allowed_slides: set[int] | None = None,
    allowed_finding_ids: set[str] | None = None,
    overflow_tolerance_pt: float = 4.0,
) -> dict:
    allowed = allowed_slides or set()
    allowed_ids = allowed_finding_ids or set()
    prs = Presentation(deck)
    overlap_findings: list[dict] = []
    overflow_findings: list[dict] = []
    unmapped_findings: list[dict] = []
    unsupported_findings: list[dict] = []
    rendered_fonts: set[str] = set()

    with fitz.open(pdf) as document:
        if len(document) != len(prs.slides):
            raise ValueError(
                "Rendered PDF slide count differs from PPTX while checking overlaps: "
                f"{len(document)} vs {len(prs.slides)}"
            )
        for slide_number, (slide, page) in enumerate(
            zip(prs.slides, document), 1
        ):
            shapes, unsupported = shape_records(prs, slide, page)
            spans = span_records(page)
            rendered_fonts.update(span["font"] for span in spans if span["font"])
            assigned, unmapped = assign_spans_to_shapes(spans, shapes)
            for finding in unsupported:
                finding["slide"] = slide_number
                annotate(
                    "unsupported_text_object",
                    finding,
                    allowed_finding_ids=allowed_ids,
                    slide_allowed=False,
                )
                unsupported_findings.append(finding)
            if unmapped:
                finding = {
                    "slide": slide_number,
                    "count": len(unmapped),
                    "textSamples": [span["text"][:80] for span in unmapped[:8]],
                    "fonts": sorted(
                        {span["font"] for span in unmapped if span["font"]}
                    ),
                }
                annotate(
                    "unmapped_rendered_text",
                    finding,
                    allowed_finding_ids=allowed_ids,
                    slide_allowed=False,
                )
                unmapped_findings.append(finding)
            for finding in detect_span_overlaps(assigned, shapes):
                finding["slide"] = slide_number
                annotate(
                    "rendered_text_overlap",
                    finding,
                    allowed_finding_ids=allowed_ids,
                    slide_allowed=slide_number in allowed,
                )
                overlap_findings.append(finding)
            for finding in detect_span_overflow(
                assigned,
                shapes,
                tolerance_pt=overflow_tolerance_pt,
            ):
                finding["slide"] = slide_number
                annotate(
                    "rendered_text_overflow",
                    finding,
                    allowed_finding_ids=allowed_ids,
                    slide_allowed=False,
                )
                overflow_findings.append(finding)

    unexpected = [
        finding for finding in overlap_findings if not finding["allowed"]
    ]
    return {
        "rendered_text_overlaps": overlap_findings,
        "unexpected_rendered_text_overlaps": unexpected,
        "rendered_text_overflow_candidates": overflow_findings,
        "unexpected_rendered_text_overflow_candidates": [
            finding for finding in overflow_findings if not finding["allowed"]
        ],
        "unmapped_rendered_text_findings": unmapped_findings,
        "unexpected_unmapped_rendered_text_findings": [
            finding for finding in unmapped_findings if not finding["allowed"]
        ],
        "unmapped_rendered_text_spans": sum(
            finding["count"] for finding in unmapped_findings
        ),
        "unsupported_text_objects": unsupported_findings,
        "unexpected_unsupported_text_objects": [
            finding for finding in unsupported_findings if not finding["allowed"]
        ],
        "rendered_fonts": sorted(rendered_fonts),
    }
